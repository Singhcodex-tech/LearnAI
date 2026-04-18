import os
import json
import re
import time
import uuid
import base64
from io import BytesIO

import requests
from flask import Flask, jsonify, render_template, request, send_file
from werkzeug.utils import secure_filename
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address

# ---------------------------------------------------------------------------
# App setup
# ---------------------------------------------------------------------------
app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "change-me-in-production")

CORS(app, origins=os.environ.get("ALLOWED_ORIGINS", "*").split(","))

# ── Law Assignment Generator (inlined) ───────────────────────────────────────
# All blueprint code inlined below after limiter setup — see _ASSIGNMENT section
# ─────────────────────────────────────────────────────────────────────────────

limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["30 per minute"],
    storage_uri="memory://",
)

# ---------------------------------------------------------------------------
# Groq API configuration
# ---------------------------------------------------------------------------
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
if not GROQ_API_KEY:
    raise ValueError(
        "GROQ_API_KEY environment variable is not set. "
        "Export it before starting the server."
    )

GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"
MODEL_NAME = "llama-3.3-70b-versatile"       # primary model — best quality
FALLBACK_MODEL = "llama-3.1-8b-instant"      # fallback if primary fails / rate-limited
MAX_TOPIC_LENGTH = 200
LAST_GROQ_CALL_AT = 0.0

# ---------------------------------------------------------------------------
# XP / Gamification constants
# ---------------------------------------------------------------------------
XP_PER_QUIZ_CORRECT = 10      # per correct answer in a 3-question quiz
XP_PER_SLIDE_UNDERSTOOD = 5
XP_STREAK_BONUS = 3           # bonus XP for every 5-slide streak

# ---------------------------------------------------------------------------
# In-memory session store
# {
#   session_id: {
#     topic, slides, performance: [...], created_at,
#     quiz_cache:   { slide_index: [questions] },
#     visual_cache: { slide_index: visual_dict },
#     notes:        { slide_index: "text" },
#     xp: int,
#     streak_slides: int,   # consecutive slides understood + quiz >= 0.6
#   }
# }
# ---------------------------------------------------------------------------
sessions: dict = {}

SESSION_TTL = 7200   # 2 hours


def _prune_sessions():
    now = time.time()
    stale = [sid for sid, s in sessions.items() if now - s["created_at"] > SESSION_TTL]
    for sid in stale:
        del sessions[sid]


# ---------------------------------------------------------------------------
# Learner profile helpers
# ---------------------------------------------------------------------------

def compute_learner_profile(performance: list, xp: int = 0, streak: int = 0) -> dict:
    """Aggregate session-wide performance into a compact profile dict."""
    if not performance:
        return {
            "summary": "new learner – no data yet",
            "avg_quiz_score": None,
            "understood_rate": None,
            "avg_time_seconds": None,
            "difficulty_hint": "normal",
            "weak_slides": [],
            "needs_review": [],
            "xp": xp,
            "streak_slides": streak,
        }

    avg_quiz = sum(p["quiz_score"] for p in performance) / len(performance)
    understood_rate = sum(1 for p in performance if p["understood"]) / len(performance)
    avg_time = sum(p["time_spent"] for p in performance) / len(performance)

    weak_slides = [
        p["slide_title"]
        for p in performance
        if p["quiz_score"] < 0.5 or not p["understood"]
    ]

    # Spaced-repetition: flag slides that need a second look
    needs_review = [
        {"slide_title": p["slide_title"], "slide_index": p["slide_index"]}
        for p in performance
        if p["quiz_score"] < 0.6 or not p["understood"]
    ]

    if avg_quiz < 0.45 or understood_rate < 0.4:
        difficulty_hint = "simplify"
    elif avg_quiz > 0.85 and understood_rate > 0.85:
        difficulty_hint = "advanced"
    else:
        difficulty_hint = "normal"

    summary_parts = []
    if avg_quiz < 0.5:
        summary_parts.append("struggling with quizzes")
    elif avg_quiz > 0.8:
        summary_parts.append("excelling at quizzes")
    if understood_rate < 0.5:
        summary_parts.append("frequently marking slides as unclear")
    if avg_time > 120:
        summary_parts.append("spending a lot of time per slide")
    elif avg_time < 15:
        summary_parts.append("progressing quickly through slides")
    if weak_slides:
        summary_parts.append(f"weak areas: {', '.join(weak_slides[:3])}")

    return {
        "summary": "; ".join(summary_parts) or "performing adequately",
        "avg_quiz_score": round(avg_quiz, 2),
        "understood_rate": round(understood_rate, 2),
        "avg_time_seconds": round(avg_time, 1),
        "difficulty_hint": difficulty_hint,
        "weak_slides": weak_slides,
        "needs_review": needs_review,
        "xp": xp,
        "streak_slides": streak,
    }


# ---------------------------------------------------------------------------
# LLM helpers
# ---------------------------------------------------------------------------

def _call_groq(
    prompt: str,
    max_tokens: int = 2000,
    system: str | None = None,
    use_fallback: bool = False,
) -> str | None:
    """
    Generic Groq call with automatic fallback to a smaller model.
    - Tries MODEL_NAME (llama-3.3-70b-versatile) first.
    - On 429 (rate-limit) or timeout, retries once with FALLBACK_MODEL.
    Returns raw text or None on error.
    """
    global LAST_GROQ_CALL_AT
    model = FALLBACK_MODEL if use_fallback else MODEL_NAME
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    try:
        # Throttle outbound LLM calls to reduce 429 bursts across endpoints.
        now = time.time()
        elapsed = now - LAST_GROQ_CALL_AT
        min_gap = 0.9
        if elapsed < min_gap:
            time.sleep(min_gap - elapsed)

        response = requests.post(
            GROQ_URL,
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": model,
                "messages": messages,
                "temperature": 0.2,
                "max_tokens": max_tokens,
            },
            timeout=45,   # increased from 30 — 70B is slower
        )
        LAST_GROQ_CALL_AT = time.time()

        # Rate-limit or server error → retry with fallback
        if response.status_code in (429, 503) and not use_fallback:
            print(f"Groq {response.status_code} on {model} — retrying with fallback model.")
            time.sleep(1.2)
            return _call_groq(prompt, max_tokens, system, use_fallback=True)
        if response.status_code in (429, 503) and use_fallback:
            print(f"Groq {response.status_code} on fallback model — backing off once.")
            time.sleep(2.0)
            return None

        if response.status_code == 401:
            print("ERROR: Groq API key is invalid or missing (401 Unauthorized). "
                  "Check your GROQ_API_KEY environment variable.")
            return None

        if response.status_code != 200:
            print(f"Groq API error {response.status_code}: {response.text[:500]}")
            return None

        content = response.json()["choices"][0]["message"]["content"]
        if not content or not content.strip():
            print(f"WARNING: Groq returned an empty response for model {model}.")
            return None
        return content

    except requests.exceptions.Timeout:
        if not use_fallback:
            print(f"Groq timeout on {model} — retrying with fallback model.")
            return _call_groq(prompt, max_tokens, system, use_fallback=True)
        print("ERROR: Groq fallback model also timed out.")
        return None
    except Exception as e:
        print(f"ERROR calling Groq: {e}")
        return None


def extract_json_array(text: str) -> str | None:
    """
    Extract the first top-level JSON array from text using bracket balancing.
    More robust than regex for deeply nested structures (slides with points arrays).
    """
    start = text.find('[')
    if start == -1:
        return None
    depth = 0
    in_string = False
    escape_next = False
    for i, ch in enumerate(text[start:], start):
        if escape_next:
            escape_next = False
            continue
        if ch == '\\' and in_string:
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '[':
            depth += 1
        elif ch == ']':
            depth -= 1
            if depth == 0:
                return text[start:i + 1]
    return None


def extract_json_object(text: str) -> str | None:
    """
    Extract the first top-level JSON object from text using bracket balancing.
    """
    start = text.find('{')
    if start == -1:
        return None
    depth = 0
    in_string = False
    escape_next = False
    for i, ch in enumerate(text[start:], start):
        if escape_next:
            escape_next = False
            continue
        if ch == '\\' and in_string:
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0:
                return text[start:i + 1]
    return None


def _extract_all_objects(text: str) -> list:
    """
    Scan through text and extract every complete top-level JSON object using
    bracket balancing. Handles deeply nested structures (points with sub_steps,
    worked_example, etc.) that single-level regex cannot match.
    Returns a list of successfully parsed dicts.
    """
    results = []
    pos = 0
    while pos < len(text):
        start = text.find('{', pos)
        if start == -1:
            break
        depth = 0
        in_string = False
        escape_next = False
        end = None
        for i, ch in enumerate(text[start:], start):
            if escape_next:
                escape_next = False
                continue
            if ch == '\\' and in_string:
                escape_next = True
                continue
            if ch == '"':
                in_string = not in_string
                continue
            if in_string:
                continue
            if ch == '{':
                depth += 1
            elif ch == '}':
                depth -= 1
                if depth == 0:
                    end = i
                    break
        if end is None:
            # Unclosed object — truncated response, stop here
            break
        candidate = text[start:end + 1]
        try:
            obj = json.loads(candidate)
            if isinstance(obj, dict):
                results.append(obj)
        except json.JSONDecodeError:
            pass
        pos = end + 1
    return results


def _repair_slides_json(raw_text: str, math_topic: bool) -> list:
    """
    Ask the model to repair malformed slide JSON into a valid array.
    Returns parsed list or [].
    """
    repair_prompt = f"""
You are a strict JSON repair engine.
Convert the following malformed model output into a valid JSON array of slides.
Return ONLY JSON array. No prose.

Required shape:
[
  {{
    "title": "Slide title",
    "points": [
      {{
        "text": "point text",
        "source_title": "source title",
        "source_url": "url or empty"
      }}
    ]
  }}
]

{ "Math topics may include inline_latex, inline_label, sub_steps, worked_example fields." if math_topic else "" }

Malformed input:
{raw_text[:12000]}
"""
    repaired = _call_groq(repair_prompt, max_tokens=3000, system="Return valid JSON only.")
    if not repaired:
        return []
    repaired = strip_markdown_fences(repaired)
    try:
        data = json.loads(repaired)
    except json.JSONDecodeError:
        arr = extract_json_array(repaired)
        if not arr:
            return []
        try:
            data = json.loads(arr)
        except json.JSONDecodeError:
            return []
    if isinstance(data, dict):
        data = data.get("slides") or data.get("data") or []
    return data if isinstance(data, list) else []


def strip_markdown_fences(text: str) -> str:
    """
    Remove ```json ... ``` or ``` ... ``` code fences that the 70B model
    sometimes adds despite being told not to.
    """
    text = re.sub(r'^```(?:json)?\s*', '', text.strip(), flags=re.IGNORECASE)
    text = re.sub(r'\s*```$', '', text.strip())
    return text.strip()


def validate_slides(slides: list) -> list:
    validated = []
    for idx, slide in enumerate(slides):
        if not isinstance(slide, dict):
            print(f"validate_slides: skipping slide {idx} — not a dict")
            continue
        title = str(slide.get("title", "")).strip()
        raw_points = slide.get("points", [])
        if not title:
            print(f"validate_slides: skipping slide {idx} — missing title")
            continue
        if not isinstance(raw_points, list):
            raw_points = []

        # Points can be either plain strings (non-math) or rich objects (math).
        # Normalise into a uniform list, preserving rich structure when present.
        points = []
        for p in raw_points:
            if isinstance(p, dict):
                text = str(p.get("text", "")).strip()
                if not text:
                    continue
                point_obj = {"text": text}
                # Per-point citation metadata
                source_title = str(p.get("source_title", "")).strip()
                source_url = str(p.get("source_url", "")).strip()
                if source_title:
                    point_obj["source_title"] = source_title
                if source_url:
                    point_obj["source_url"] = source_url
                # inline equation
                il = str(p.get("inline_latex", "")).strip()
                lbl = str(p.get("inline_label", "")).strip()
                if il:
                    point_obj["inline_latex"] = il
                    point_obj["inline_label"] = lbl
                # sub-steps
                ss = p.get("sub_steps", [])
                if isinstance(ss, list):
                    ss = [str(s).strip() for s in ss if str(s).strip()]
                    if ss:
                        point_obj["sub_steps"] = ss
                points.append(point_obj)
            elif isinstance(p, str):
                text = p.strip()
                if text:
                    points.append({"text": text})

        if len(points) < 1:
            print(f"validate_slides: skipping slide '{title}' — no valid points parsed")
            continue

        entry: dict = {"title": title, "points": points[:6]}

        # Preserve worked example (math slides)
        worked = slide.get("worked_example")
        if isinstance(worked, dict):
            problem = str(worked.get("problem", "")).strip()
            steps = worked.get("steps", [])
            answer = str(worked.get("answer", "")).strip()
            if isinstance(steps, list):
                steps = [str(s).strip() for s in steps if str(s).strip()]
            if problem and steps:
                entry["worked_example"] = {
                    "problem": problem,
                    "steps": steps,
                    "answer": answer,
                }

        validated.append(entry)
    return validated


def _expand_short_points_with_api(slides: list, min_words: int = 50) -> list:
    """
    API-only expansion pass: if in-depth points are too short, ask the model
    to rewrite only those points to at least `min_words`.
    """
    short_exists = False
    for s in slides:
        for p in s.get("points", []):
            if isinstance(p, dict) and len(str(p.get("text", "")).split()) < min_words:
                short_exists = True
                break
        if short_exists:
            break
    if not short_exists:
        return slides

    prompt = f"""
You are editing slide JSON.
Expand only short point texts so each point has at least {min_words} words.
Keep meaning intact. Keep title, source_title, source_url, and any math fields.
Return ONLY valid JSON array, no prose.

Input slides:
{json.dumps(slides, ensure_ascii=False)}
"""
    raw = _call_groq(prompt, max_tokens=3200, system="Return valid JSON only.")
    if not raw:
        return slides
    raw = strip_markdown_fences(raw)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        arr = extract_json_array(raw)
        if not arr:
            return slides
        try:
            data = json.loads(arr)
        except json.JSONDecodeError:
            return slides
    if isinstance(data, dict):
        data = data.get("slides") or data.get("data") or []
    if not isinstance(data, list):
        return slides
    fixed = validate_slides(data)
    return fixed or slides


def _normalize_point_word_lengths(slides: list, explanation_mode: str) -> list:
    """
    Fix #3: Audit-only pass — never truncate point text.

    The old implementation truncated at max_w words, destroying the detail
    that was explicitly requested in the prompt (45-55 words in_depth).
    This version only logs short points so we can monitor quality without
    destroying any content.
    """
    if explanation_mode == "brief":
        min_w = 15
    else:
        min_w = 35

    for slide in slides:
        for p in slide.get("points", []):
            if not isinstance(p, dict):
                continue
            text = str(p.get("text", "")).strip()
            if not text:
                continue
            word_count = len(text.split())
            if word_count < min_w:
                print(
                    f"_normalize_point_word_lengths: short point ({word_count} words) "
                    f"in slide '{slide.get('title', '?')}' — keeping as-is."
                )
            # Never modify p["text"] — return slides untouched
    return slides


def generate_slides_rescue(topic: str, explanation_mode: str = "in_depth") -> list:
    """
    Fix #4: Rescue generation using the same 1-slide-per-call + role approach as
    generate_slides(), but forced onto the fallback model (llama-3.1-8b-instant).
    Iterates all 12 _DECK_ROLES so the rescue can always fill a complete deck,
    unlike the old version that was capped at 5 slides / 3 points.
    """
    point_length_rule = (
        "- Every point text should be around 50 words (target range 45-55 words)"
        if explanation_mode == "in_depth"
        else "- Every point text should be around 20 to 25 words"
    )
    depth_note = (
        "\nNOTE: Explain in depth with concrete details."
        if explanation_mode == "in_depth"
        else "\nNOTE: Keep explanations brief and easy to scan."
    )
    math_topic = is_math_topic(topic)

    rescued: list = []
    for slide_idx, (role_title, role_instruction) in enumerate(_DECK_ROLES):
        slide_number = slide_idx + 1
        slide = _generate_single_slide(
            topic=topic,
            slide_number=slide_number,
            role_title=role_title,
            role_instruction=role_instruction,
            math_topic=math_topic,
            difficulty_hint="",
            depth_note=depth_note,
            point_length_rule=point_length_rule,
            retry=True,   # force fallback model
        )
        if slide is not None:
            rescued.append(slide)
        print(f"generate_slides_rescue: {len(rescued)}/{slide_number} slides rescued.")

    return rescued


# ---------------------------------------------------------------------------
# Structured Learning Control — Subject / Mode / Depth prompt builder
# ---------------------------------------------------------------------------

VALID_MODES  = {"Learn", "Practice", "Test"}
VALID_DEPTHS = {"Beginner", "Exam", "Advanced"}

# Every subject gets its own mode × depth matrix.
# Keys must exactly match what the frontend sends (title-cased).
_SUBJECT_MODE_INSTRUCTIONS: dict[str, dict[str, str]] = {
    # ── STEM ────────────────────────────────────────────────────────────
    "Math": {
        "Learn":    "Explain the concept step-by-step with clear worked examples. "
                    "Show each calculation stage and define every symbol used.",
        "Practice": "Generate realistic practice problems (numeric, word-problem, proof) "
                    "with full solutions and reasoning shown at each step.",
        "Test":     "Generate exam-style questions ONLY. Do NOT include answers, "
                    "hints, or solutions anywhere in the slide content.",
    },
    "Physics": {
        "Learn":    "Explain the physical law or phenomenon with intuition first, "
                    "then derive the governing equations step-by-step with SI units.",
        "Practice": "Provide quantitative physics problems with full worked solutions, "
                    "including free-body diagrams or vector descriptions where relevant.",
        "Test":     "Pose physics problems with given data and ask for derivations or "
                    "numerical answers ONLY. Do NOT include solutions or hints.",
    },
    "Chemistry": {
        "Learn":    "Explain chemical concepts with reaction mechanisms, bonding diagrams, "
                    "and real-world examples of compounds or processes.",
        "Practice": "Provide stoichiometry, reaction-balancing, or mechanism exercises "
                    "with complete solutions and mole/unit tracking shown.",
        "Test":     "Present chemical problems or name-the-product tasks ONLY. "
                    "Do NOT include answers, structures, or explanations.",
    },
    "Biology": {
        "Learn":    "Explain biological processes and structures using clear diagrams "
                    "described in text, real organism examples, and function-structure links.",
        "Practice": "Provide labelling exercises, process-ordering tasks, or case-study "
                    "questions with model answers and common-mistake notes.",
        "Test":     "Generate structured biology questions and diagrams-to-label ONLY. "
                    "Do NOT include answers or model responses.",
    },
    "Coding": {
        "Learn":    "Explain the programming concept clearly, then show a concise working "
                    "code example. Break down each line and explain why it is written that way.",
        "Practice": "Provide a coding challenge with input/output examples and a partial "
                    "skeleton or hint. Include the full correct solution with explanation.",
        "Test":     "Present programming problems and specifications ONLY. Do NOT include "
                    "solutions, pseudocode answers, or implementation hints.",
    },
    "Computer Science": {
        "Learn":    "Explain CS theory (algorithms, data structures, complexity) with "
                    "step-by-step traces, pseudocode, and Big-O analysis.",
        "Practice": "Provide algorithm design or analysis exercises with full solutions, "
                    "including time/space complexity justification.",
        "Test":     "Pose CS theory or algorithm questions ONLY. Do NOT include "
                    "solutions, proofs, or complexity answers.",
    },
    "Statistics": {
        "Learn":    "Explain statistical concepts with formulas, real datasets, and "
                    "interpretation of results in plain language.",
        "Practice": "Provide calculation exercises (hypothesis tests, confidence intervals, "
                    "regression) with full worked solutions and interpretation.",
        "Test":     "Present statistical problems with data tables ONLY. "
                    "Do NOT include calculations, p-values, or conclusions.",
    },
    # ── HUMANITIES ──────────────────────────────────────────────────────
    "English": {
        "Learn":    "Explain grammar rules or literary concepts with clear examples "
                    "and counter-examples. Use real sentences to illustrate every rule.",
        "Practice": "Provide writing or grammar correction exercises. Include model "
                    "answers and explanations of common mistakes.",
        "Test":     "Generate comprehension tasks, essay prompts, or grammar exercises "
                    "ONLY. Do NOT include model answers or corrections.",
    },
    "History": {
        "Learn":    "Narrate historical events with causes, key figures, dates, and "
                    "long-term consequences. Connect events to broader patterns.",
        "Practice": "Provide source-analysis or essay-structure exercises based on "
                    "real historical documents, with model responses.",
        "Test":     "Pose source-based or essay questions ONLY. "
                    "Do NOT include model answers or mark schemes.",
    },
    "Geography": {
        "Learn":    "Explain physical or human geography concepts with real place examples, "
                    "statistics, and cause-effect chains.",
        "Practice": "Provide case-study analysis or map-interpretation exercises "
                    "with full model answers and data commentary.",
        "Test":     "Present geography questions or data-response tasks ONLY. "
                    "Do NOT include answers or explanations.",
    },
    "Philosophy": {
        "Learn":    "Introduce the philosophical argument, key thinkers, objections, "
                    "and counter-arguments with clear logical structure.",
        "Practice": "Provide argument-analysis or essay-plan tasks with model responses "
                    "demonstrating critical evaluation techniques.",
        "Test":     "Pose philosophical essay questions or argument-analysis tasks ONLY. "
                    "Do NOT include model answers or evaluation commentary.",
    },
    # ── APPLIED / PROFESSIONAL ──────────────────────────────────────────
    "Economics": {
        "Learn":    "Explain economic theory with real-world data, diagrams described "
                    "in text (supply/demand, curves), and policy implications.",
        "Practice": "Provide calculation or diagram-analysis exercises (elasticity, "
                    "GDP, market equilibrium) with full solutions.",
        "Test":     "Present economics data-response or essay questions ONLY. "
                    "Do NOT include answers, calculations, or diagram descriptions.",
    },
    "Business": {
        "Learn":    "Explain business concepts (strategy, finance, marketing, HR) with "
                    "real company examples, frameworks, and key metrics.",
        "Practice": "Provide case-study analysis tasks or financial calculation exercises "
                    "with model answers and evaluation criteria.",
        "Test":     "Pose business case-study or essay questions ONLY. "
                    "Do NOT include model answers or mark schemes.",
    },
    "Law": {
        "Learn":    "Explain legal principles, landmark cases, and statutes clearly. "
                    "Use the IRAC framework (Issue, Rule, Application, Conclusion).",
        "Practice": "Provide problem scenarios requiring legal analysis with model "
                    "IRAC answers and note on jurisdiction where relevant.",
        "Test":     "Present legal problem questions or essay prompts ONLY. "
                    "Do NOT include model answers, case citations, or analysis.",
    },
    "Medicine": {
        "Learn":    "Explain anatomy, physiology, or clinical concepts with mechanisms, "
                    "pathophysiology, and real clinical presentations.",
        "Practice": "Provide clinical vignette exercises with full diagnostic reasoning, "
                    "differential diagnoses, and management plans.",
        "Test":     "Present clinical vignettes or MCQ stems ONLY. "
                    "Do NOT include diagnoses, reasoning, or management answers.",
    },
    "Psychology": {
        "Learn":    "Explain psychological theories, key studies, and real-world "
                    "applications with evaluation of strengths and limitations.",
        "Practice": "Provide essay-plan or study-evaluation exercises with model "
                    "answers demonstrating PEEL structure.",
        "Test":     "Pose psychology essay or short-answer questions ONLY. "
                    "Do NOT include model answers or study details.",
    },
}

VALID_SUBJECTS = set(_SUBJECT_MODE_INSTRUCTIONS.keys())

# Depth modifiers apply to ALL subjects uniformly
_DEPTH_INSTRUCTIONS: dict[str, str] = {
    "Beginner":  "Use simple, jargon-free language. Assume no prior knowledge. "
                 "Include everyday analogies and small, accessible examples.",
    "Exam":      "Use structured, moderately detailed explanations suited to exam "
                 "preparation. Balance clarity with academic rigour and use precise terminology.",
    "Advanced":  "Use precise technical language and provide deep, challenging content. "
                 "Include edge cases, nuances, competing theories, and expert-level detail.",
}


def build_learning_context(subject: str, mode: str, depth: str) -> str:
    """
    Return an instruction block prepended to every slide prompt when the user
    picks Subject + Mode + Depth.  Returns "" if any value is missing/invalid
    so the default university-professor prompt is used unchanged.
    All three controls must be set for the optimised prompt to activate.
    """
    subject = (subject or "").strip()
    mode    = (mode    or "").strip().title()
    depth   = (depth   or "").strip().title()

    # Normalise subject capitalisation against known keys
    subject_match = next(
        (s for s in VALID_SUBJECTS if s.lower() == subject.lower()), ""
    )

    if not subject_match or mode not in VALID_MODES or depth not in VALID_DEPTHS:
        # Incomplete selection → fall through to default professor prompt (no change)
        return ""

    mode_text  = _SUBJECT_MODE_INSTRUCTIONS[subject_match][mode]
    depth_text = _DEPTH_INSTRUCTIONS[depth]

    return (
        f"LEARNING CONTEXT — Subject: {subject_match} | Mode: {mode} | Depth: {depth}\n"
        f"Mode instruction: {mode_text}\n"
        f"Depth instruction: {depth_text}\n"
    )


# ---------------------------------------------------------------------------
# Math topic detection
# ---------------------------------------------------------------------------

MATH_KEYWORDS = {
    # Pure math branches
    "algebra", "calculus", "geometry", "trigonometry", "statistics",
    "probability", "number theory", "linear algebra", "differential equations",
    "integral", "derivative", "matrix", "vector", "polynomial", "logarithm",
    "exponential", "binomial", "permutation", "combination", "fourier",
    "laplace", "topology", "set theory", "graph theory", "arithmetic",
    "fraction", "equation", "inequality", "quadratic", "sequence", "series",
    "limit", "continuity", "complex number", "real analysis", "discrete math",
    # Applied math / physics math
    "optimization", "eigenvalue", "gradient", "divergence", "curl",
    "bayes", "hypothesis testing", "regression", "variance", "standard deviation",
    "normal distribution", "correlation", "integration", "differentiation",
}

def is_math_topic(topic: str) -> bool:
    """Return True if the topic is mathematical in nature."""
    lower = topic.lower()
    return any(kw in lower for kw in MATH_KEYWORDS)


# ---------------------------------------------------------------------------
# Slide generation — 1 slide per API call (Fix #1 + #2)
# ---------------------------------------------------------------------------

TOTAL_SLIDES = 12   # target number of slides per session

# Fix #2 — explicit role for every slide position in the 12-slide deck.
# Each role tells the model exactly what job this slide must do, so every
# slide is distinct and purposeful (no more generic, repetitive output).
_DECK_ROLES: list[tuple[str, str]] = [
    ("Introduction & Definition",
     "Define the topic clearly. Explain what it is, where it comes from, and why it matters. "
     "Include the core vocabulary a learner must know before going further."),
    ("Historical Background & Origin",
     "Cover the history, origin, and key milestones of the topic. Who developed it, when, and why? "
     "Mention specific names, dates, or events."),
    ("Core Principles & Mechanisms",
     "Explain the fundamental principles, rules, or laws that govern the topic. "
     "Focus on HOW it works at a mechanistic level."),
    ("Key Components & Structure",
     "Break down the major parts, components, layers, or categories. "
     "Use precise terminology and explain what each component does."),
    ("Types, Variants & Classification",
     "Describe the main types, variants, or sub-categories of the topic. "
     "Explain the distinguishing features and when each variant is used."),
    ("Step-by-Step Process or Workflow",
     "Walk through the main process, algorithm, lifecycle, or workflow step by step. "
     "Be concrete — use numbered or ordered logic."),
    ("Real-World Applications & Use Cases",
     "Give specific, concrete real-world examples and applications. "
     "Mention industries, companies, products, or specific scenarios with names and numbers."),
    ("Advantages, Strengths & Benefits",
     "Detail the key advantages, strengths, and benefits. "
     "Support each claim with a reason or a quantitative fact where possible."),
    ("Limitations, Challenges & Criticisms",
     "Explain the key limitations, drawbacks, trade-offs, and criticisms. "
     "Be specific — vague negatives are not useful."),
    ("Comparison with Related Concepts",
     "Compare and contrast the topic with closely related concepts, alternatives, or competing approaches. "
     "Highlight what makes each distinct and when to prefer one over the other."),
    ("Current Trends, Research & Future Directions",
     "Cover the latest developments, active research areas, and where the field is heading. "
     "Mention specific technologies, papers, or events from recent years."),
    ("Summary, Key Takeaways & Next Steps",
     "Synthesise the entire topic into the most important insights a learner should remember. "
     "End with practical next steps or resources for deeper study."),
]


def _build_single_slide_prompt(
    topic: str,
    slide_number: int,          # 1-based position in the deck
    role_title: str,
    role_instruction: str,
    math_topic: bool,
    difficulty_hint: str,
    depth_note: str,
    point_length_rule: str,
    learning_context: str = "",
) -> tuple[str, int]:
    """
    Build a prompt that requests exactly ONE slide.
    Single-slide requests produce tiny JSON payloads (~300-500 tokens),
    making JSON truncation and corruption impossible.
    Returns (prompt_str, max_tokens).
    """
    position_ctx = (
        f"You are building slide {slide_number} of {TOTAL_SLIDES} "
        f"in a deck about: {topic}\n"
        f"This slide's role: {role_title}\n"
        f"What this slide must cover: {role_instruction}\n"
    )

    # Prepend structured learning context if provided
    ctx_block = f"{learning_context}\n" if learning_context else ""

    if math_topic:
        prompt = (
            "You are an expert mathematics professor creating one rigorous, exam-quality slide.\n\n"
            f"{ctx_block}"
            f"{position_ctx}"
            f"{difficulty_hint}{depth_note}\n\n"
            "Return ONLY a valid JSON object (a single slide). No markdown fences. No extra text.\n\n"
            "Required format:\n"
            "{\n"
            "  \"title\": \"Slide title that reflects the role\",\n"
            "  \"points\": [\n"
            "    {\n"
            "      \"text\": \"One sentence: state the rule or theorem precisely. No filler words.\",\n"
            "      \"source_title\": \"Real credible source (e.g. Stewart Calculus, MIT OCW, Wikipedia)\",\n"
            "      \"source_url\": \"Direct URL or empty string\",\n"
            "      \"inline_latex\": \"\\\\frac{-b \\\\pm \\\\sqrt{b^2-4ac}}{2a}\",\n"
            "      \"inline_label\": \"Formula name (e.g. Quadratic Formula)\",\n"
            "      \"sub_steps\": [\n"
            "        \"Step 1 — Identify: $$a = ...,\\\\; b = ...,\\\\; c = ...$$\",\n"
            "        \"Step 2 — Substitute: $$x = \\\\frac{-b \\\\pm \\\\sqrt{b^2-4ac}}{2a}$$\",\n"
            "        \"Step 3 — Simplify: $$x = ...$$\",\n"
            "        \"Step 4 — Verify: $$...$$\"\n"
            "      ]\n"
            "    }\n"
            "  ],\n"
            "  \"worked_example\": {\n"
            "    \"problem\": \"State the problem using LaTeX, e.g. \\\"Solve $$x^2 - 5x + 6 = 0$$\\\"\",\n"
            "    \"steps\": [\n"
            "      \"Step 1 — Setup: $$a=1,\\\\; b=-5,\\\\; c=6$$\",\n"
            "      \"Step 2 — Substitute: $$x = \\\\frac{5 \\\\pm \\\\sqrt{25-24}}{2}$$\",\n"
            "      \"Step 3 — Simplify: $$x = \\\\frac{5 \\\\pm 1}{2}$$\",\n"
            "      \"Step 4 — Solve: $$x = 3 \\\\text{ or } x = 2$$\"\n"
            "    ],\n"
            "    \"answer\": \"$$x = 2$$ or $$x = 3$$\"\n"
            "  }\n"
            "}\n\n"
            "STRICT RULES:\n"
            "- Exactly 4 points\n"
            f"{point_length_rule}\n"
            "- Every point MUST be an object with: text, source_title, source_url, inline_latex, inline_label, sub_steps\n"
            "- text: ONE sentence — state the mathematical rule, theorem, or definition precisely. "
            "  NO long explanations. May contain inline $...$ LaTeX.\n"
            "- inline_latex: the key formula as a display LaTeX expression (double backslashes: \\\\frac, \\\\sqrt, \\\\int, \\\\pm)\n"
            "- inline_label: short name for the formula (e.g. 'Chain Rule', 'Euler's Formula')\n"
            "- sub_steps: exactly 4 steps, EACH formatted as 'Step N — <Action>: $$<LaTeX expression>$$'. "
            "  ZERO prose sentences — every step is a concrete mathematical operation in LaTeX. "
            "  Valid actions: Identify, Set up, Substitute, Expand, Factor, Differentiate, Integrate, "
            "  Simplify, Solve, Apply, Verify, Evaluate, Rearrange.\n"
            "- worked_example.problem: state clearly with LaTeX math\n"
            "- worked_example.steps: 4–5 steps each 'Step N — <Action>: $$<LaTeX>$$' — "
            "  every single intermediate result MUST be in $$...$$\n"
            "- worked_example.answer: final result in LaTeX e.g. '$$x = 2$$ or $$x = 3$$'\n"
            "- Do NOT use bullet points, prose explanations, or long sentences in sub_steps\n"
            "- Do NOT repeat formulas or examples from other slides\n"
            "- Output ONLY the JSON object — no array wrapper, no prose\n"
        )
        max_tok = 1600
    else:
        prompt = (
            "You are an expert university professor creating one deeply detailed, lecture-quality slide.\n\n"
            f"{ctx_block}"
            f"{position_ctx}"
            f"{difficulty_hint}{depth_note}\n\n"
            "Return ONLY a valid JSON object (a single slide). No markdown fences. No extra text.\n\n"
            "Required format:\n"
            "{\n"
            "  \"title\": \"Slide title that reflects the role\",\n"
            "  \"points\": [\n"
            "    {\n"
            "      \"text\": \"Detailed, specific explanation sentence\",\n"
            "      \"source_title\": \"Real credible source title\",\n"
            "      \"source_url\": \"Direct URL or empty string\"\n"
            "    }\n"
            "  ]\n"
            "}\n\n"
            "STRICT RULES:\n"
            "- Exactly 4 bullet points\n"
            f"{point_length_rule}\n"
            "- Every point MUST be an object with keys: text, source_title, source_url\n"
            "- Every point must explain the 'what' and briefly the 'why'\n"
            "- NO vague lines like 'It is important' or 'Has many applications'\n"
            "- Every point MUST include at least one of: a definition, a mechanism, a real example "
            "with a name or number, a comparison, a cause-effect relationship, or a formula\n"
            "- Do NOT repeat information from other slides\n"
            "- Output ONLY the JSON object — no array wrapper, no prose\n"
        )
        max_tok = 1200

    return prompt, max_tok


def _generate_single_slide(
    topic: str,
    slide_number: int,
    role_title: str,
    role_instruction: str,
    math_topic: bool,
    difficulty_hint: str,
    depth_note: str,
    point_length_rule: str,
    retry: bool = False,
    learning_context: str = "",
) -> dict | None:
    """
    Call the API for exactly one slide and return a validated slide dict (or None).
    """
    prompt, max_tok = _build_single_slide_prompt(
        topic, slide_number, role_title, role_instruction,
        math_topic, difficulty_hint, depth_note, point_length_rule,
        learning_context=learning_context,
    )

    raw = _call_groq(prompt, max_tokens=max_tok, use_fallback=retry)
    if not raw:
        if not retry:
            print(f"_generate_single_slide [{slide_number}]: API returned None, retrying with fallback.")
            return _generate_single_slide(
                topic, slide_number, role_title, role_instruction,
                math_topic, difficulty_hint, depth_note, point_length_rule,
                retry=True, learning_context=learning_context,
            )
        print(f"_generate_single_slide [{slide_number}]: giving up.")
        return None

    raw = strip_markdown_fences(raw)
    print(f"_generate_single_slide [{slide_number}]: preview {raw[:120]!r}")

    # Parse — expect a single JSON object
    obj = None
    try:
        obj = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"_generate_single_slide [{slide_number}]: json.loads failed ({e}), trying extraction.")
        # Maybe model wrapped it in an array
        arr_str = extract_json_array(raw)
        if arr_str:
            try:
                arr = json.loads(arr_str)
                if isinstance(arr, list) and arr:
                    obj = arr[0]
            except json.JSONDecodeError:
                pass
        if obj is None:
            obj_str = extract_json_object(raw)
            if obj_str:
                try:
                    obj = json.loads(obj_str)
                except json.JSONDecodeError:
                    pass
        if obj is None:
            objs = _extract_all_objects(raw)
            if objs:
                obj = objs[0]

    if not isinstance(obj, dict):
        if not retry:
            return _generate_single_slide(
                topic, slide_number, role_title, role_instruction,
                math_topic, difficulty_hint, depth_note, point_length_rule,
                retry=True, learning_context=learning_context,
            )
        return None

    validated = validate_slides([obj])
    if not validated:
        if not retry:
            return _generate_single_slide(
                topic, slide_number, role_title, role_instruction,
                math_topic, difficulty_hint, depth_note, point_length_rule,
                retry=True, learning_context=learning_context,
            )
        return None

    print(f"_generate_single_slide [{slide_number}]: OK — '{validated[0]['title']}'")
    return validated[0]


def generate_slides(
    topic: str,
    learner_profile: dict | None = None,
    explanation_mode: str = "in_depth",
    compact_mode: bool = False,
    retry: bool = False,
    learning_context: str = "",
) -> list:
    explanation_mode = str(explanation_mode or "in_depth").strip().lower()
    if explanation_mode not in {"brief", "in_depth"}:
        explanation_mode = "in_depth"

    if explanation_mode == "brief":
        depth_note = (
            "\nNOTE: Keep explanations brief and easy to scan. "
            "Use concise wording and only essential detail."
        )
        point_length_rule = "- Every point text should be around 20 to 25 words"
    else:
        depth_note = (
            "\nNOTE: Explain in depth. Include richer context, clear reasoning, and concrete details."
        )
        point_length_rule = "- Every point text should be around 50 words (target range 45–55 words)"

    difficulty_hint = ""
    if learner_profile:
        hint = learner_profile.get("difficulty_hint", "normal")
        profile_summary = learner_profile.get("summary", "")
        if hint == "simplify":
            difficulty_hint = (
                "\nNOTE: This learner has been struggling. Use simpler language, "
                "more analogies, and concrete step-by-step examples. Avoid jargon. "
                f"Learner profile: {profile_summary}\n"
            )
        elif hint == "advanced":
            difficulty_hint = (
                "\nNOTE: This learner is advanced. Include more technical depth, "
                "edge cases, and challenging concepts. "
                f"Learner profile: {profile_summary}\n"
            )

    math_topic = is_math_topic(topic)

    # ── 1-slide-per-call generation ──────────────────────────────────────────
    # Each API call requests exactly 1 slide using an explicit role from
    # _DECK_ROLES. This guarantees:
    #   • JSON output is tiny (~300–500 tokens) — impossible to truncate/corrupt
    #   • Each slide has a distinct, purposeful job — no generic repetition
    all_slides: list = []
    for slide_idx, (role_title, role_instruction) in enumerate(_DECK_ROLES):
        slide_number = slide_idx + 1
        slide = _generate_single_slide(
            topic=topic,
            slide_number=slide_number,
            role_title=role_title,
            role_instruction=role_instruction,
            math_topic=math_topic,
            difficulty_hint=difficulty_hint,
            depth_note=depth_note,
            point_length_rule=point_length_rule,
            learning_context=learning_context,
        )
        if slide is not None:
            all_slides.append(slide)
        else:
            print(f"generate_slides: slide {slide_number} failed — skipping.")

        print(f"generate_slides: {len(all_slides)}/{slide_number} slides so far.")

    # ── Rescue: fill any missing slides via fallback model ───────────────────
    if len(all_slides) < TOTAL_SLIDES:
        missing = TOTAL_SLIDES - len(all_slides)
        print(f"generate_slides: {missing} slides missing — running rescue pass.")
        rescue = generate_slides_rescue(topic, explanation_mode)
        needed = TOTAL_SLIDES - len(all_slides)
        all_slides.extend(rescue[:needed])

    # Fix #3 — do NOT truncate points; audit only (normalize is now a no-op truncation-wise)
    if all_slides:
        all_slides = _normalize_point_word_lengths(all_slides, explanation_mode)

    print(f"generate_slides: final count = {len(all_slides)} slides.")
    return all_slides


# ---------------------------------------------------------------------------
# Quiz generation
# ---------------------------------------------------------------------------

def generate_quiz_for_slide(slide: dict, learner_profile: dict | None = None, retry: bool = False) -> list:
    """Generate 3 adaptive MCQ questions for a slide."""
    difficulty_note = ""
    if learner_profile:
        hint = learner_profile.get("difficulty_hint", "normal")
        if hint == "simplify":
            difficulty_note = "Make questions straightforward and conceptual. Avoid trick questions."
        elif hint == "advanced":
            difficulty_note = "Make questions challenging. Include application-based and analytical questions."

    # Detect math slides by presence of inline_latex in any point
    is_math_slide = any(
        isinstance(p, dict) and p.get("inline_latex")
        for p in slide.get("points", [])
    )
    math_note = (
        "This is a mathematics slide. "
        "Questions and options may contain LaTeX expressions using $...$ for inline math "
        "and $$...$$ for display equations. "
        "Include at least one question that requires evaluating or identifying a formula. "
        "Options involving numbers or expressions must use LaTeX."
    ) if is_math_slide else ""

    prompt = f"""
You are a university professor writing a quiz based on this slide.

Slide Title: {slide['title']}
Slide Content:
{json.dumps(slide['points'], indent=2)}

{difficulty_note}
{math_note}

Generate exactly 3 multiple-choice questions that test understanding of this slide's content.

Return ONLY a valid JSON array. No extra text, no markdown.

Format:
[
  {{
    "question": "Clear, specific question? (use $LaTeX$ for any math expressions)",
    "options": ["Option A (use $LaTeX$ for math)", "Option B", "Option C", "Option D"],
    "correct": 0,
    "explanation": "Concise explanation referencing the slide. Use $LaTeX$ for any math."
  }}
]

RULES:
- All 4 options must be plausible (no obviously wrong options)
- "correct" is the 0-based index of the correct option
- Questions must be directly answerable from the slide content
- Vary question types: recall, application, conceptual
- Keep options concise (under 20 words each)
- For math slides: at least one question must test formula recall or application
- Do NOT output anything outside JSON
"""

    raw_text = _call_groq(prompt, max_tokens=1800)  # increased for 70B verbosity
    if not raw_text:
        return []

    raw_text = strip_markdown_fences(raw_text)

    try:
        questions = json.loads(raw_text)
    except json.JSONDecodeError:
        json_part = extract_json_array(raw_text)
        if not json_part:
            if not retry:
                return generate_quiz_for_slide(slide, learner_profile, retry=True)
            return []
        try:
            questions = json.loads(json_part)
        except json.JSONDecodeError:
            return []

    if not isinstance(questions, list):
        return []

    validated = []
    for q in questions:
        if not isinstance(q, dict):
            continue
        if not all(k in q for k in ("question", "options", "correct", "explanation")):
            continue
        if not isinstance(q["options"], list) or len(q["options"]) != 4:
            continue
        if not isinstance(q["correct"], int) or not (0 <= q["correct"] <= 3):
            continue
        validated.append({
            "question": str(q["question"]).strip(),
            "options": [str(o).strip() for o in q["options"]],
            "correct": q["correct"],
            "explanation": str(q["explanation"]).strip(),
        })

    return validated[:3]


# ---------------------------------------------------------------------------
# Re-teaching (simplified slide regeneration)
# ---------------------------------------------------------------------------

def reteach_slide(slide: dict, retry: bool = False) -> dict | None:
    """Regenerate a slide with simpler language, more examples, and analogies.
    For math slides (those with equations/worked_example), also regenerates
    a simpler worked example with more detailed steps."""
    is_math = "equations" in slide or "worked_example" in slide

    if is_math:
        points_block = json.dumps(slide.get("points", []), indent=2)
        worked_block = json.dumps(slide.get("worked_example", {}), indent=2)
        reteach_prompt = (
            "A student did not understand this mathematics slide.\n"
            "Rewrite it with SIMPLER language and a NEW, EASIER worked example.\n\n"
            "Original Slide Title: %s\n"
            "Original Points:\n%s\n\n"
            "Original Worked Example:\n%s\n\n"
            "Rules:\n"
            "- Keep the same title\n"
            "- Rewrite each point as a rich object with: text, source_title, source_url, inline_latex, inline_label, sub_steps\n"
            "- Use simpler language and everyday analogies\n"
            "- inline_latex: same formula but write a clearer inline_label and simpler sub_steps\n"
            "- sub_steps: 3-5 steps, each following \'Step N — <Label>: <explanation with LaTeX>\' format; "
            "every number or expression must use $...$ or $$...$$ LaTeX delimiters\n"
            "- worked_example: a NEW simpler numeric problem, 4-6 steps with more arithmetic detail; "
            "every step must follow \'Step N — <Label>: <LaTeX working>\' and use $$...$$ for calculations\n"
            "- Keep 4 to 6 points\n\n"
            "Return ONLY a valid JSON object. No extra text.\n\n"
            "Format:\n"
            "{\n"
            "  \"title\": \"Same title\",\n"
            "  \"points\": [\n"
            "    {\n"
            "      \"text\": \"Simpler explanation sentence\",\n"
            "      \"source_title\": \"Credible source title for this bullet\",\n"
            "      \"source_url\": \"Direct URL/DOI for this bullet source, or empty string if unavailable\",\n"
            "      \"inline_latex\": \"\\\\frac{a}{b}\",\n"
            "      \"inline_label\": \"Formula name\",\n"
            "      \"sub_steps\": [\"Step 1 — reason: detail\", \"Step 2 — reason: detail\"]\n"
            "    }\n"
            "  ],\n"
            "  \"worked_example\": {\n"
            "    \"problem\": \"A simpler problem\",\n"
            "    \"steps\": [\"Step 1 — Setup: ...\", \"Step 2 — Compute: ...\"],\n"
            "    \"answer\": \"Final answer\"\n"
            "  },\n"
            "  \"reteach\": true\n"
            "}\n"
        ) % (slide["title"], points_block, worked_block)
        prompt = reteach_prompt
        max_tokens = 2000
    else:
        prompt = f"""
A student did not understand this slide. Rewrite it to be clearer and easier to understand.

Original Slide Title: {slide['title']}
Original Content:
{json.dumps(slide['points'], indent=2)}

Rewrite rules:
- Keep the same title
- Use simpler language — explain as if to a high-school student
- Add a real-world analogy or everyday example for each concept
- Break complex ideas into step-by-step explanations
- Replace technical jargon with plain English where possible
- Each point must still be a complete, informative sentence
- Keep 4 to 6 bullet points

Return ONLY a valid JSON object. No extra text.

Format:
{{
  "title": "Same title",
  "points": [
    "Simplified explanation with example",
    "Another clearer explanation"
  ],
  "reteach": true
}}
"""
        max_tokens = 1000

    raw_text = _call_groq(prompt, max_tokens=max_tokens)
    if not raw_text:
        return None

    raw_text = strip_markdown_fences(raw_text)

    try:
        slide_data = json.loads(raw_text)
    except json.JSONDecodeError:
        json_part = extract_json_object(raw_text)
        if not json_part:
            if not retry:
                return reteach_slide(slide, retry=True)
            return None
        try:
            slide_data = json.loads(json_part)
        except json.JSONDecodeError:
            return None

    if not isinstance(slide_data, dict):
        return None

    title = str(slide_data.get("title", slide["title"])).strip()
    points = slide_data.get("points", [])
    if not isinstance(points, list):
        return None
    points = [str(p).strip() for p in points if str(p).strip()]
    if len(points) < 2:
        return None

    # Re-parse points using the same normalisation as validate_slides
    raw_points = slide_data.get("points", [])
    if not isinstance(raw_points, list):
        raw_points = []
    norm_points = []
    for p in raw_points:
        if isinstance(p, dict):
            text = str(p.get("text", "")).strip()
            if not text:
                continue
            po = {"text": text}
            source_title = str(p.get("source_title", "")).strip()
            source_url = str(p.get("source_url", "")).strip()
            if source_title:
                po["source_title"] = source_title
            if source_url:
                po["source_url"] = source_url
            il = str(p.get("inline_latex", "")).strip()
            lbl = str(p.get("inline_label", "")).strip()
            if il:
                po["inline_latex"] = il
                po["inline_label"] = lbl
            ss = p.get("sub_steps", [])
            if isinstance(ss, list):
                ss = [str(s).strip() for s in ss if str(s).strip()]
                if ss:
                    po["sub_steps"] = ss
            norm_points.append(po)
        elif isinstance(p, str):
            t = p.strip()
            if t:
                norm_points.append({"text": t})

    if len(norm_points) < 2:
        return None

    result: dict = {"title": title, "points": norm_points[:6], "reteach": True}

    if is_math:
        worked = slide_data.get("worked_example")
        if isinstance(worked, dict):
            problem = str(worked.get("problem", "")).strip()
            steps = worked.get("steps", [])
            answer = str(worked.get("answer", "")).strip()
            if isinstance(steps, list):
                steps = [str(s).strip() for s in steps if str(s).strip()]
            if problem and steps:
                result["worked_example"] = {
                    "problem": problem,
                    "steps": steps,
                    "answer": answer,
                }

    return result


# ---------------------------------------------------------------------------
# NEW: Topic summary generation
# ---------------------------------------------------------------------------

def generate_topic_summary(topic: str, slides: list, performance: list) -> str | None:
    """
    Generate a concise end-of-session summary with key takeaways.
    Personalises the summary based on weak slides if performance data exists.
    For math topics, includes key formulas in LaTeX.
    """
    slide_titles = [s["title"] for s in slides]
    weak_titles = [p["slide_title"] for p in performance if p["quiz_score"] < 0.6 or not p["understood"]]

    weak_note = ""
    if weak_titles:
        weak_note = f"\nThe learner struggled with: {', '.join(weak_titles[:5])}. Briefly flag those areas for extra study."

    math_summary_note = ""
    if is_math_topic(topic):
        # Collect any inline_latex from slides to give the LLM context
        formulas = []
        for s in slides:
            for p in s.get("points", []):
                if isinstance(p, dict) and p.get("inline_latex"):
                    lbl = p.get("inline_label", "")
                    lat = p["inline_latex"]
                    formulas.append(f"{lbl}: $${lat}$$" if lbl else f"$${lat}$$")
        if formulas:
            math_summary_note = (
                "\nKey formulas covered (include these in your summary using LaTeX $...$ notation):\n"
                + "\n".join(formulas[:8])
            )
        math_summary_note += (
            "\nUse LaTeX for ALL mathematical expressions in your summary: "
            "inline with $expression$ and display equations with $$expression$$."
        )

    prompt = f"""
You are an expert tutor. A student just finished a learning session on "{topic}".

The session covered these slides:
{json.dumps(slide_titles, indent=2)}
{weak_note}
{math_summary_note}

Write a clear, motivating end-of-session summary for the student. Include:
1. A 2-3 sentence overview of what was learned
2. 5-7 key takeaways as concise bullet points (include key formulas in LaTeX for math topics)
3. A short "What to explore next" section with 2-3 related topics
4. An encouraging closing line

Keep the tone friendly and academic. Plain text — no JSON, no markdown headers.
"""

    return _call_groq(prompt, max_tokens=900)


# ---------------------------------------------------------------------------
# NEW: AI Tutor chat (contextual Q&A on current slide)
# ---------------------------------------------------------------------------

# Extended keyword set for question-level math detection (broader than
# is_math_topic which targets topic names; this targets student questions).
_MATH_QUESTION_KEYWORDS = MATH_KEYWORDS | {
    "solve", "calculate", "compute", "prove", "simplify", "differentiate",
    "integrate", "find", "evaluate", "expand", "factor", "factorise",
    "determinant", "eigenvalue", "eigenvector", "sum", "product",
    "formula", "expression", "simplify", "roots", "zeros",
}


def is_mathematical_question(question: str) -> bool:
    """Return True if the student's question is asking for a mathematical solution."""
    lower = question.lower()
    return any(kw in lower for kw in _MATH_QUESTION_KEYWORDS)


def generate_math_solution(question: str, slide_points: list) -> str | None:
    """
    Return a structured, step-by-step mathematical solution.
    Every step is labelled 'Step N — <action>:' with full LaTeX.
    Answer is on a final 'Step N — Answer:' line.
    """
    context = ""
    if slide_points:
        context_lines = [
            (p["text"] if isinstance(p, dict) else str(p))
            for p in slide_points[:3]
        ]
        context = "Relevant slide context:\n" + "\n".join(context_lines) + "\n\n"

    system = (
        "You are a precise mathematics tutor. "
        "Solve the problem completely, showing every intermediate calculation. "
        "Format your entire response as clearly numbered steps — no bullet points, "
        "no dashes, no markdown headers or symbols. "
        "Each step must be on its own line labelled exactly "
        "\'Step N — <action label>: <explanation>\' where N is the step number. "
        "Use LaTeX for ALL mathematical expressions: "
        "inline with $expression$ and display equations on their own line with $$expression$$. "
        "Every number, variable, formula, and equation must be wrapped in LaTeX. "
        "Explain what you are doing and why at each step. "
        "End with a line labelled \'Step N — Answer: $$result$$\' containing the final result "
        "with units where applicable. "
        "Never use bullet points, dashes, or plain text for mathematical expressions."
    )

    prompt = (
        f"{context}"
        f"Solve the following problem completely, showing every step:\n\n"
        f"{question}\n\n"
        "Use this exact format for every step:\n"
        "Step 1 — <action label>: <explanation with LaTeX e.g. Substituting $a=2$ into $$x = \\frac{-b}{2a}$$>\n"
        "Step 2 — <action label>: ...\n"
        "...\n"
        "Step N — Answer: $$<final result>$$"
    )

    raw = _call_groq(prompt, max_tokens=1200, system=system)
    if raw:
        raw = strip_markdown_fences(raw)
    return raw


def _topic_is_math_question(topic: str) -> bool:
    """
    Return True when the user's topic input looks like a concrete math problem
    to solve rather than a concept to study (e.g. "solve x^2+5x+6=0" vs "quadratic equations").
    Combines action-verb detection with numeric / operator pattern matching.
    """
    import re as _re
    lower = topic.lower().strip()
    # Action verbs that signal a problem to solve
    action_verbs = {
        "solve", "calculate", "compute", "simplify", "differentiate",
        "integrate", "find", "evaluate", "expand", "factor", "factorise",
        "prove", "derive", "verify", "determine", "show that",
    }
    has_verb = any(lower.startswith(v) or f" {v} " in lower for v in action_verbs)
    # Numeric / operator patterns: digits with operators, equals sign, fractions, etc.
    has_math_pattern = bool(_re.search(
        r'(\d[\d\s]*[+\-*/^=<>])|'   # number followed by operator
        r'([a-z]\^?\d)|'              # variable like x^2, x2
        r'(=\s*\d)|'                  # = some number
        r'[∫∑∏√±÷×]|'               # math symbols
        r'\\frac|\\sqrt|\\int|\\sum', # LaTeX fragments
        topic
    ))
    return has_verb or has_math_pattern


def _math_solution_to_slide(question: str, solution_text: str) -> dict:
    """
    Convert raw step-by-step solution text (from generate_math_solution) into
    a single slide dict compatible with the frontend renderer.

    The solution text uses the format:
        Step 1 — <action>: <explanation with LaTeX>
        ...
        Step N — Answer: $$result$$

    We store all steps as sub_steps inside a single point so the existing
    math-slide renderer displays them with LaTeX rendering — no bullet points.
    """
    lines = [l.strip() for l in solution_text.strip().splitlines() if l.strip()]

    steps = []
    answer_line = ""
    for line in lines:
        if line.lower().startswith("step") and "answer" in line.lower():
            answer_line = line
            steps.append(line)
        elif line.lower().startswith("step"):
            steps.append(line)

    # Build the slide as a single rich point with all steps as sub_steps
    point = {
        "text": f"Step-by-step solution for: {question}",
        "source_title": "",
        "source_url": "",
        "inline_latex": "",
        "inline_label": "Solution",
        "sub_steps": steps if steps else lines,
    }

    slide: dict = {
        "title": f"Solution: {question[:80]}",
        "points": [point],
        "is_solution_slide": True,
    }

    # Extract answer for worked_example.answer display
    if answer_line:
        slide["worked_example"] = {
            "problem": question,
            "steps": steps[:-1] if len(steps) > 1 else steps,
            "answer": answer_line.split(":", 1)[-1].strip() if ":" in answer_line else answer_line,
        }

    return slide


def answer_student_question(question: str, slide: dict, topic: str, profile: dict) -> str | None:
    # Mathematical questions bypass the slide-context prompt entirely and go
    # straight to a dedicated solver that returns structured step-by-step output.
    if is_mathematical_question(question):
        return generate_math_solution(question, slide.get("points", []))

    difficulty = profile.get("difficulty_hint", "normal")
    tone_note = {
        "simplify": "Use simple, jargon-free language with everyday analogies.",
        "advanced": "Use precise technical language and include deeper details.",
        "normal": "Use clear, academic language suitable for an undergraduate student.",
    }.get(difficulty, "")

    # Detect if the current slide is math-heavy so tutor uses LaTeX
    slide_is_math = any(
        isinstance(p, dict) and p.get("inline_latex")
        for p in slide.get("points", [])
    ) or is_math_topic(topic)
    latex_note = (
        "Use LaTeX for ALL mathematical expressions in your answer: "
        "inline with $expression$ and display equations with $$expression$$. "
        "Never write equations or formulas in plain text."
    ) if slide_is_math else ""

    system = (
        f"You are a friendly, knowledgeable university tutor teaching '{topic}'. "
        f"{tone_note} {latex_note} "
        "Answer only questions related to the course material. "
        "If the question is off-topic, gently redirect the student. "
        "Keep your answer focused and under 200 words."
    )

    prompt = f"""Current slide the student is on:
Title: {slide['title']}
Content:
{json.dumps(slide['points'], indent=2)}

Student's question: {question}
"""

    return _call_groq(prompt, max_tokens=600, system=system)


# ---------------------------------------------------------------------------
# Visual generation
# ---------------------------------------------------------------------------

def generate_visual_for_slide(slide: dict, retry: bool = False) -> dict | None:
    """
    Decide the best visual type for a slide and return structured data
    that the frontend can render as an SVG/HTML diagram.
    """
    prompt = f"""Return ONLY valid JSON object (no markdown).
Create one educational diagram for this slide.

Slide title: {slide['title']}
Slide points:
{json.dumps(slide.get('points', [])[:4], ensure_ascii=False)}

Use one type from: flowchart, cycle, comparison, timeline, pyramid, mindmap.
Keep structure compact (4 to 6 elements).

Output format:
{{
  "type": "mindmap",
  "title": "Short title",
  "data": {{
    "center": "Main idea",
    "branches": [
      {{"label": "Branch label", "items": ["item one", "item two"]}}
    ]
  }}
}}
If you choose another type, keep the same top-level keys and valid data for that type.
"""

    raw = _call_groq(prompt, max_tokens=650, use_fallback=True)
    if not raw:
        return None

    raw = strip_markdown_fences(raw)

    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        obj = extract_json_object(raw)
        if not obj:
            if not retry:
                return generate_visual_for_slide(slide, retry=True)
            return None
        try:
            data = json.loads(obj)
        except json.JSONDecodeError:
            return None

    if not isinstance(data, dict):
        return None
    if data.get("type") not in {"flowchart", "cycle", "comparison", "timeline", "pyramid", "mindmap"}:
        return None
    if not isinstance(data.get("data"), dict):
        return None

    return data


def build_fallback_visual(slide: dict) -> dict:
    """
    Deterministic fallback visual so frontend always receives a renderable diagram.
    """
    title = str(slide.get("title", "Slide Concept")).strip() or "Slide Concept"
    raw_points = slide.get("points", [])
    texts = []
    for p in raw_points:
        t = str(p.get("text", "") if isinstance(p, dict) else p).strip()
        if t:
            texts.append(t)
    if not texts:
        texts = [title]

    branches = []
    for i, t in enumerate(texts[:5]):
        words = t.split()
        label = " ".join(words[:4]) or f"Idea {i + 1}"
        items = []
        if len(words) > 4:
            items.append(" ".join(words[4:10]))
        if len(words) > 10:
            items.append(" ".join(words[10:16]))
        if not items:
            items = ["Key point from this slide", "Related concept to review"]
        branches.append({"label": label[:48], "items": [it[:72] for it in items[:2]]})

    return {
        "type": "mindmap",
        "title": title[:48],
        "data": {
            "center": title[:48],
            "branches": branches[:5],
        },
    }


# ---------------------------------------------------------------------------
# Sources generation
# ---------------------------------------------------------------------------

def generate_sources_for_slide(slide: dict, topic: str, retry: bool = False) -> list:
    """
    Generate a list of credible reference sources for the content on a given slide.
    Returns a list of source dicts: { title, type, description, url_hint }
    """
    points_text = "\n".join(
        (p["text"] if isinstance(p, dict) else str(p))
        for p in slide.get("points", [])
    )

    prompt = f"""You are an academic librarian. For the slide below, list the 3 to 5 most relevant, credible reference sources a student could use to learn more.

Topic: {topic}
Slide Title: {slide['title']}
Slide Content:
{points_text}

Return ONLY a valid JSON array. No extra text.

Format:
[
  {{
    "title": "Full name of the source (book, paper, website, or course)",
    "type": "book" | "paper" | "website" | "course" | "encyclopedia",
    "authors": "Author(s) or organization (e.g. 'Goodfellow et al.' or 'Wikipedia')",
    "year": "Publication year or 'ongoing' for websites",
    "description": "One sentence explaining what this source covers and why it is relevant to this slide.",
    "url_hint": "A plausible URL or DOI (e.g. 'https://en.wikipedia.org/wiki/...' or 'https://arxiv.org/...'). Use an empty string if unknown."
  }}
]

RULES:
- Only cite real, well-known sources (textbooks, Wikipedia, Khan Academy, arXiv, MDN, official docs, etc.)
- Do NOT invent fictional papers or fake DOIs
- Match sources to the exact content of the slide (not just the general topic)
- Prefer freely accessible sources where possible
- Do NOT output anything outside the JSON array
"""

    raw = _call_groq(prompt, max_tokens=1200)
    if not raw:
        return []

    raw = strip_markdown_fences(raw)

    try:
        sources = json.loads(raw)
    except json.JSONDecodeError:
        arr = extract_json_array(raw)
        if not arr:
            if not retry:
                return generate_sources_for_slide(slide, topic, retry=True)
            return []
        try:
            sources = json.loads(arr)
        except json.JSONDecodeError:
            return []

    if not isinstance(sources, list):
        return []

    validated = []
    for s in sources:
        if not isinstance(s, dict):
            continue
        title = str(s.get("title", "")).strip()
        if not title:
            continue
        validated.append({
            "title": title,
            "type": str(s.get("type", "website")).strip(),
            "authors": str(s.get("authors", "")).strip(),
            "year": str(s.get("year", "")).strip(),
            "description": str(s.get("description", "")).strip(),
            "url_hint": str(s.get("url_hint", "")).strip(),
        })

    return validated[:5]


# ---------------------------------------------------------------------------
# XP / Streak helpers
# ---------------------------------------------------------------------------

def _award_xp(session: dict, quiz_score: float, understood: bool) -> dict:
    """
    Update XP and streak in-place and return a dict with the delta info.
    """
    correct_answers = round(quiz_score * 3)  # 3-question quiz
    earned = correct_answers * XP_PER_QUIZ_CORRECT
    if understood:
        earned += XP_PER_SLIDE_UNDERSTOOD

    # Streak logic
    if quiz_score >= 0.6 and understood:
        session["streak_slides"] = session.get("streak_slides", 0) + 1
    else:
        session["streak_slides"] = 0

    # Streak bonus every 5 slides
    streak = session["streak_slides"]
    if streak > 0 and streak % 5 == 0:
        earned += XP_STREAK_BONUS

    session["xp"] = session.get("xp", 0) + earned

    return {
        "xp_earned": earned,
        "total_xp": session["xp"],
        "streak_slides": session["streak_slides"],
    }


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route("/")
def home():
    return render_template("index.html")


@app.route("/health")
@app.route("/health.php")
def health():
    """Simple health check for deployment monitoring."""
    return jsonify({"status": "ok", "model": MODEL_NAME, "sessions": len(sessions)})


# ---------------------------------------------------------------------------
# File Upload — PDF & Image → extract topic text, then feed into /generate
# ---------------------------------------------------------------------------
#
# Llama-3.3-70b (Groq) is TEXT-ONLY — it cannot accept binary files.
# Solution:
#   • PDF  → extract text with PyMuPDF (fitz)  → send extracted text to Llama
#   • Image → send base64 to llava-v1.5-7b-4096-preview (free vision model on Groq)
#             → get text description → send description to Llama as topic
#
# Install deps once:  pip install PyMuPDF
#

LLAVA_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"  # free vision model on Groq
UPLOAD_MAX_MB = 10
ALLOWED_EXTENSIONS = {"pdf", "png", "jpg", "jpeg", "webp", "gif"}


def _allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def _extract_pdf_text(file_bytes: bytes, max_chars: int = 6000) -> str:
    """Extract plain text from PDF bytes using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError(
            "PyMuPDF is not installed. Run: pip install PyMuPDF"
        )
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    full_text = "\n".join(text_parts).strip()
    return full_text[:max_chars]


def _describe_image_with_llava(image_bytes: bytes, media_type: str) -> str | None:
    """
    Send image to llava-v1.5-7b (free vision model on Groq) and get a description.
    Returns plain-text description or None on error.
    """
    b64 = base64.b64encode(image_bytes).decode("utf-8")
    payload = {
        "model": LLAVA_MODEL,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{media_type};base64,{b64}"
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "Describe this image in detail. Extract all visible text, "
                            "labels, diagrams, equations, and key concepts. "
                            "Focus on educational content that could be used as a learning topic."
                        ),
                    },
                ],
            }
        ],
        "max_tokens": 800,
        "temperature": 0.2,
    }
    try:
        resp = requests.post(
            GROQ_URL,
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=45,
        )
        if resp.status_code != 200:
            print(f"LLaVA error {resp.status_code}: {resp.text[:300]}")
            return None
        return resp.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        print(f"LLaVA call failed: {e}")
        return None


@app.route("/upload_file", methods=["POST"])
@app.route("/upload_file.php", methods=["POST"])
@limiter.limit("5 per minute")
def upload_file():
    """
    Accept a PDF or image upload. Extract its content and return a suggested topic
    string plus a content preview — the frontend then calls /generate with that topic.

    Form fields:
      file  — the uploaded file (multipart/form-data)

    Returns:
      {
        "topic":   "Suggested topic / title for the session",
        "preview": "First ~500 chars of extracted content",
        "type":    "pdf" | "image"
      }
    """
    if "file" not in request.files:
        return jsonify({"error": "No file part in request"}), 400

    f = request.files["file"]
    if not f or not f.filename:
        return jsonify({"error": "No file selected"}), 400
    if not _allowed_file(f.filename):
        return jsonify({
            "error": f"Unsupported file type. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        }), 400

    file_bytes = f.read()
    if len(file_bytes) > UPLOAD_MAX_MB * 1024 * 1024:
        return jsonify({"error": f"File too large. Max {UPLOAD_MAX_MB} MB."}), 413

    ext = f.filename.rsplit(".", 1)[1].lower()
    filename_safe = secure_filename(f.filename)

    # ── PDF branch ────────────────────────────────────────────────────────────
    if ext == "pdf":
        try:
            extracted = _extract_pdf_text(file_bytes)
        except RuntimeError as e:
            return jsonify({"error": str(e)}), 500

        if not extracted:
            return jsonify({"error": "Could not extract text from PDF. It may be a scanned image-only PDF."}), 422

        # Ask Llama to infer a concise topic title from the extracted text
        title_prompt = (
            f"Read the following text extracted from a PDF and reply with ONLY a short, "
            f"specific topic title (5–10 words) that best describes what this document is about. "
            f"No quotes, no explanation — just the topic title.\n\n{extracted[:2000]}"
        )
        topic = _call_groq(title_prompt, max_tokens=60) or filename_safe
        topic = topic.strip().strip('"').strip("'")

        return jsonify({
            "topic":   topic[:MAX_TOPIC_LENGTH],
            "preview": extracted[:500],
            "type":    "pdf",
            "full_text": extracted,   # frontend can pass this back as context if needed
        })

    # ── Image branch ─────────────────────────────────────────────────────────
    else:
        mime_map = {
            "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "png": "image/png", "webp": "image/webp", "gif": "image/gif",
        }
        media_type = mime_map.get(ext, "image/jpeg")

        description = _describe_image_with_llava(file_bytes, media_type)
        if not description:
            return jsonify({"error": "Could not analyse image. Try a clearer image or PDF."}), 422

        # Ask Llama to distil description into a topic title
        title_prompt = (
            f"Based on the following image description, reply with ONLY a short, "
            f"specific topic title (5–10 words) suitable for a learning session. "
            f"No quotes, no explanation.\n\n{description[:1500]}"
        )
        topic = _call_groq(title_prompt, max_tokens=60) or filename_safe
        topic = topic.strip().strip('"').strip("'")

        return jsonify({
            "topic":   topic[:MAX_TOPIC_LENGTH],
            "preview": description[:500],
            "type":    "image",
        })


@app.route("/generate", methods=["POST"])
@app.route("/generate.php", methods=["POST"])
@limiter.limit("10 per minute")
def generate():
    """
    Start a new learning session.
    Body: {
      "topic": "...",
      "explanation_mode": "brief" | "in_depth",
      "subject": "Math" | "English" | "Coding",   (optional)
      "mode":    "Learn" | "Practice" | "Test",    (optional)
      "depth":   "Beginner" | "Exam" | "Advanced"  (optional)
    }
    Returns: { "session_id": "...", "slides": [...] }
    """
    _prune_sessions()
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid or missing JSON body"}), 400

    topic = str(data.get("topic", "")).strip()
    explanation_mode = str(data.get("explanation_mode", "in_depth")).strip().lower()
    if not topic:
        return jsonify({"error": "No topic provided"}), 400
    if len(topic) > MAX_TOPIC_LENGTH:
        return jsonify({"error": f"Topic must be {MAX_TOPIC_LENGTH} characters or fewer"}), 400
    if explanation_mode not in {"brief", "in_depth"}:
        explanation_mode = "in_depth"

    # Structured learning controls (optional)
    subject = str(data.get("subject", "")).strip()
    mode    = str(data.get("mode",    "")).strip()
    depth   = str(data.get("depth",   "")).strip()
    learning_context = build_learning_context(subject, mode, depth)

    # ── Math-question fast-path ──────────────────────────────────────────────
    # If the topic looks like a concrete math problem (e.g. "solve x^2+5x+6=0"),
    # skip 12-slide theory generation entirely and return ONE solution slide.
    if _topic_is_math_question(topic):
        print(f"[/generate] Math question detected — generating single solution slide for: {topic!r}")
        solution_text = generate_math_solution(topic, [])
        if solution_text:
            slides = [_math_solution_to_slide(topic, solution_text)]
            session_id = str(uuid.uuid4())
            sessions[session_id] = {
                "topic": topic,
                "slides": slides,
                "performance": [],
                "created_at": time.time(),
                "quiz_cache": {},
                "visual_cache": {},
                "sources_cache": {},
                "notes": {},
                "explanation_mode": explanation_mode,
                "subject": subject,
                "mode": mode,
                "depth": depth,
                "xp": 0,
                "streak_slides": 0,
                "is_solution_session": True,
            }
            return jsonify({
                "session_id": session_id,
                "slides": slides,
                "explanation_mode": explanation_mode,
                "is_solution_session": True,
            })
        # Fall through to normal generation if solution fails
        print(f"[/generate] Math solution generation failed — falling back to normal slide generation.")

    slides = []
    for attempt in range(3):
        slides = generate_slides(topic, explanation_mode=explanation_mode,
                                 compact_mode=False, learning_context=learning_context)
        if slides:
            break
        print(f"[/generate] Attempt {attempt + 1}/3 failed for topic='{topic}'. Retrying...")
        time.sleep(1.2 * (attempt + 1))
    # If normal-size generation still fails, automatically switch to compact payload.
    if not slides:
        print(f"[/generate] Switching to compact generation for topic='{topic}'.")
        for attempt in range(3):
            slides = generate_slides(topic, explanation_mode=explanation_mode,
                                     compact_mode=True, learning_context=learning_context)
            if slides:
                break
            print(f"[/generate] Compact attempt {attempt + 1}/3 failed for topic='{topic}'. Retrying...")
            time.sleep(1.5 * (attempt + 1))
    if not slides:
        print(f"[/generate] Standard generation failed for topic='{topic}'. Trying rescue generation.")
        slides = generate_slides_rescue(topic, explanation_mode=explanation_mode)
    if not slides:
        print(f"[/generate] No slides produced for topic='{topic}' after retries.")
        return jsonify({
            "error": "Slide generation failed after retries. Please try once more."
        }), 500

    session_id = str(uuid.uuid4())
    sessions[session_id] = {
        "topic": topic,
        "slides": slides,
        "performance": [],
        "created_at": time.time(),
        "quiz_cache": {},
        "visual_cache": {},
        "sources_cache": {},
        "notes": {},
        "explanation_mode": explanation_mode,
        "subject": subject,
        "mode": mode,
        "depth": depth,
        "xp": 0,
        "streak_slides": 0,
        "is_solution_session": False,
    }

    return jsonify({"session_id": session_id, "slides": slides, "explanation_mode": explanation_mode, "is_solution_session": False})


@app.route("/quiz", methods=["POST"])
@app.route("/quiz.php", methods=["POST"])
@limiter.limit("20 per minute")
def quiz():
    """
    Generate adaptive quiz for a specific slide (cached per slide).
    Body: { "session_id": "...", "slide_index": 0 }
    Returns: { "questions": [...] }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    slide_index = data.get("slide_index", 0)

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    slides = session["slides"]
    if not (0 <= slide_index < len(slides)):
        return jsonify({"error": "Invalid slide index"}), 400

    cache_key = str(slide_index)
    if cache_key in session["quiz_cache"]:
        return jsonify({"questions": session["quiz_cache"][cache_key], "cached": True})

    profile = compute_learner_profile(session["performance"], session.get("xp", 0), session.get("streak_slides", 0))
    questions = generate_quiz_for_slide(slides[slide_index], profile)

    if not questions:
        return jsonify({"error": "Failed to generate quiz. Please try again."}), 500

    session["quiz_cache"][cache_key] = questions
    return jsonify({"questions": questions})


@app.route("/feedback", methods=["POST"])
@app.route("/feedback.php", methods=["POST"])
@limiter.limit("30 per minute")
def feedback():
    """
    Record learner performance for a slide and award XP.
    Body: {
        "session_id": "...",
        "slide_index": 0,
        "slide_title": "...",
        "time_spent": 45,
        "understood": true,
        "quiz_score": 0.67,
        "quiz_answers": [0, 2, 1]
    }
    Returns: { "profile": {...}, "message": "...", "xp": {...} }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    record = {
        "slide_index": int(data.get("slide_index", 0)),
        "slide_title": str(data.get("slide_title", "")).strip(),
        "time_spent": float(data.get("time_spent", 0)),
        "understood": bool(data.get("understood", True)),
        "quiz_score": float(data.get("quiz_score", 1.0)),
        "quiz_answers": data.get("quiz_answers", []),
        "timestamp": time.time(),
    }

    session["performance"] = [
        p for p in session["performance"]
        if p["slide_index"] != record["slide_index"]
    ]
    session["performance"].append(record)

    # Award XP
    xp_result = _award_xp(session, record["quiz_score"], record["understood"])

    profile = compute_learner_profile(session["performance"], session.get("xp", 0), session.get("streak_slides", 0))

    if record["quiz_score"] >= 0.8 and record["understood"]:
        message = "Excellent work! You have a strong grasp of this topic."
    elif record["quiz_score"] >= 0.5 or record["understood"]:
        message = "Good effort. Review any points you found confusing before moving on."
    else:
        message = "This topic needs more attention. Try the simplified explanation."

    return jsonify({"profile": profile, "message": message, "xp": xp_result})


@app.route("/reteach", methods=["POST"])
@app.route("/reteach.php", methods=["POST"])
@limiter.limit("10 per minute")
def reteach():
    """
    Return a simplified version of a slide.
    Body: { "session_id": "...", "slide_index": 0 }
    Returns: { "slide": { "title": "...", "points": [...], "reteach": true } }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    slide_index = data.get("slide_index", 0)

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    slides = session["slides"]
    if not (0 <= slide_index < len(slides)):
        return jsonify({"error": "Invalid slide index"}), 400

    simplified = reteach_slide(slides[slide_index])
    if not simplified:
        return jsonify({"error": "Failed to generate simplified slide."}), 500

    return jsonify({"slide": simplified})


def _stats_payload(session_id: str):
    """
    Return full session analytics.
    Returns: { "topic", "total_slides", "completed_slides", "profile", "performance" }
    """
    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    profile = compute_learner_profile(session["performance"], session.get("xp", 0), session.get("streak_slides", 0))

    return jsonify({
        "topic": session["topic"],
        "total_slides": len(session["slides"]),
        "completed_slides": len(session["performance"]),
        "profile": profile,
        "performance": session["performance"],
        "created_at": session["created_at"],
    })


@app.route("/stats/<session_id>", methods=["GET"])
def stats(session_id):
    return _stats_payload(session_id)


@app.route("/stats.php", methods=["GET"])
def stats_php():
    session_id = str(request.args.get("session_id", "")).strip()
    if not session_id:
        return jsonify({"error": "Missing session_id"}), 400
    return _stats_payload(session_id)


@app.route("/visual", methods=["POST"])
@app.route("/visual.php", methods=["POST"])
@limiter.limit("15 per minute")
def visual():
    """
    Generate a visual diagram for a specific slide (cached per slide).
    Body:  { "session_id": "...", "slide_index": 0 }
    Returns: { "visual": { "type": "...", "title": "...", "data": {...} } }
    """
    req_data = request.get_json(silent=True)
    if not req_data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = req_data.get("session_id", "")
    slide_index = req_data.get("slide_index", 0)

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    slides = session["slides"]
    if not (0 <= slide_index < len(slides)):
        return jsonify({"error": "Invalid slide index"}), 400

    cache_key = str(slide_index)
    if cache_key in session["visual_cache"]:
        return jsonify({"visual": session["visual_cache"][cache_key], "cached": True})

    result = generate_visual_for_slide(slides[slide_index])
    if not result:
        time.sleep(0.8)
        result = generate_visual_for_slide(slides[slide_index], retry=True)
    if not result:
        result = build_fallback_visual(slides[slide_index])

    session["visual_cache"][cache_key] = result
    return jsonify({"visual": result})


# ---------------------------------------------------------------------------
# NEW: AI Tutor chat endpoint
# ---------------------------------------------------------------------------

@app.route("/chat", methods=["POST"])
@app.route("/chat.php", methods=["POST"])
@limiter.limit("20 per minute")
def chat():
    """
    Ask the AI tutor a question about the current slide.
    Body: {
        "session_id": "...",
        "slide_index": 0,
        "question": "Why does X happen?"
    }
    Returns: { "answer": "..." }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    slide_index = data.get("slide_index", 0)
    question = str(data.get("question", "")).strip()

    if not question:
        return jsonify({"error": "No question provided"}), 400
    if len(question) > 500:
        return jsonify({"error": "Question too long (max 500 characters)"}), 400

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    slides = session["slides"]
    if not (0 <= slide_index < len(slides)):
        return jsonify({"error": "Invalid slide index"}), 400

    profile = compute_learner_profile(session["performance"], session.get("xp", 0), session.get("streak_slides", 0))
    answer = answer_student_question(question, slides[slide_index], session["topic"], profile)

    if not answer:
        return jsonify({"error": "Tutor is unavailable. Please try again."}), 500

    return jsonify({"answer": answer.strip()})


# ---------------------------------------------------------------------------
# Math Tutor — helpers
# ---------------------------------------------------------------------------

DIRECT_SOLVE_TRIGGERS = {
    "solve", "simplify", "calculate", "compute", "evaluate", "integrate",
    "differentiate", "derive", "factor", "expand", "find", "prove",
    "what is", "what's", "how much", "how many",
}

def _is_direct_solve(question: str) -> bool:
    """Return True if the user wants a full solution right now."""
    q = question.lower()
    return any(kw in q for kw in DIRECT_SOLVE_TRIGGERS)


def _math_direct_solve(question: str, difficulty: str) -> str | None:
    diff_instruction = {
        "Beginner":  "Show every micro-step. Write each arithmetic operation on its own line.",
        "Exam":      "Show all working. Every intermediate result must appear as a display equation.",
        "Advanced":  "Use rigorous notation, state domain restrictions, and justify each transformation.",
    }.get(difficulty, "Show every intermediate result as a display equation.")

    system = (
        "You are a mathematics solution engine. "
        "Your ONLY job is to output a clean, structured, step-by-step mathematical solution. "
        "STRICT OUTPUT RULES — violating any rule is an error:\n"
        "1. Every step MUST be on its own line in EXACTLY this format: "
        "   'Step N — <Short Action Label>: $$<LaTeX expression>$$'\n"
        "2. The action label must be a short math verb: e.g. 'Expand', 'Factor', 'Substitute', "
        "   'Differentiate', 'Simplify', 'Integrate', 'Apply quadratic formula', 'Solve for x'.\n"
        "3. EVERY mathematical expression — equations, formulas, numbers, variables — MUST be "
        "   in LaTeX using $...$ (inline) or $$...$$ (display). NO plain-text math.\n"
        "4. The final step MUST be: 'Step N — Answer: $$<result>$$'\n"
        "5. ZERO prose, ZERO explanatory sentences, ZERO bullet points, ZERO markdown headers.\n"
        "6. If a brief clarification of what was done is needed, append it in parentheses "
        "   AFTER the LaTeX on the SAME line only — maximum 6 words.\n"
        f"7. {diff_instruction}\n"
        "Output ONLY the numbered steps. Nothing before Step 1. Nothing after the Answer step."
    )
    prompt = (
        f"Solve completely:\n\n{question}\n\n"
        "Format each line EXACTLY as:\n"
        "Step 1 — <Action>: $$<LaTeX>$$\n"
        "Step 2 — <Action>: $$<LaTeX>$$\n"
        "...\n"
        "Step N — Answer: $$<final result>$$\n\n"
        "Every equation must be in LaTeX. No prose. No bullets. No headers."
    )
    raw = _call_groq(prompt, max_tokens=1800, system=system)
    if raw:
        raw = strip_markdown_fences(raw)
    return raw


def _generate_math_problem(topic: str, difficulty: str) -> dict | None:
    """
    Ask the model to produce a structured math problem JSON:
    { problem, steps[], final_answer, similar_problems[] }
    """
    diff_note = {
        "Beginner": "appropriate for a beginner (arithmetic / basic algebra)",
        "Exam":     "exam-level difficulty (multi-step, real techniques required)",
        "Advanced": "advanced (proof-based or multi-concept)",
    }.get(difficulty, "intermediate difficulty")

    step_count = {"Beginner": "3–4", "Exam": "4–6", "Advanced": "5–7"}.get(difficulty, "4–5")

    system = (
        "You are a math tutor. Return ONLY a valid JSON object. "
        "No markdown fences. No extra text."
    )
    prompt = f"""Generate a math problem on the topic: "{topic}".
Difficulty: {diff_note}.

Return ONLY this JSON shape:
{{
  "problem": "Full problem statement. Use LaTeX for ALL math: inline $expression$ or display $$expression$$. Every number, variable, and formula must be in LaTeX.",
  "steps": [
    "Step 1 — Setup: identify given values $a = ...$, $b = ...$",
    "Step 2 — Apply: substitute into $$formula$$",
    "Step 3 — Compute: $$intermediate = result$$",
    "Step N — Answer: $$final = value$$"
  ],
  "final_answer": "$$answer$$ with units where applicable",
  "similar_problems": [
    "Similar problem 1 (use LaTeX for any math)",
    "Similar problem 2"
  ]
}}

Rules:
- steps[] must be {step_count} items, each a complete verifiable calculation step
- Every step must follow the format "Step N — <Label>: <working with LaTeX>"
- Every mathematical expression in steps, problem, and final_answer MUST use LaTeX ($...$ or $$...$$)
- final_answer must be precise, correct, and in LaTeX
- similar_problems must be 1–2 items at the same difficulty
- Use correct math — verify every calculation before writing
- Do NOT use plain text for any equation, variable, or number"""

    raw = _call_groq(prompt, max_tokens=1200, system=system)
    if not raw:
        return None
    raw = strip_markdown_fences(raw)
    try:
        obj = json.loads(raw)
    except json.JSONDecodeError:
        s = extract_json_object(raw)
        if not s:
            return None
        try:
            obj = json.loads(s)
        except json.JSONDecodeError:
            return None

    if not isinstance(obj, dict):
        return None
    if not obj.get("problem") or not isinstance(obj.get("steps"), list):
        return None
    return obj


def _validate_math_step(user_answer: str, expected_step: str, problem: str, topic: str) -> dict:
    """
    Return { correct: bool, hint: str|None }
    Asks the model to judge correctness tolerantly (accept equivalent forms).
    Hint uses LaTeX so the frontend can render it properly.
    """
    system = "You are a strict but fair math evaluator. Return ONLY valid JSON. No prose outside JSON."
    prompt = f"""Problem: {problem}
Expected step: {expected_step}
Student's answer: {user_answer}

Is the student's answer mathematically equivalent to or correctly performing the expected step?
Be tolerant of notation differences (e.g. "x=2" vs "$x = 2$") but strict about mathematical correctness.

If incorrect, write a hint that:
1. Points out the specific error in one sentence
2. Shows the correct intermediate expression using LaTeX (e.g. "Try $$x = \\frac{{-b}}{{2a}}$$")
3. Never gives away the full answer

Return ONLY valid JSON:
{{"correct": true/false, "hint": "hint with LaTeX if incorrect, null if correct"}}"""

    raw = _call_groq(prompt, max_tokens=300, system=system)
    if not raw:
        return {"correct": False, "hint": "Could not validate. Please try again."}
    raw = strip_markdown_fences(raw)
    try:
        obj = json.loads(raw)
        return {
            "correct": bool(obj.get("correct", False)),
            "hint": obj.get("hint") or None,
        }
    except Exception:
        s = extract_json_object(raw)
        if s:
            try:
                obj = json.loads(s)
                return {"correct": bool(obj.get("correct", False)), "hint": obj.get("hint")}
            except Exception:
                pass
        return {"correct": False, "hint": "Could not validate. Please try again."}


# ---------------------------------------------------------------------------
# Math Tutor — /math_tutor endpoint
# ---------------------------------------------------------------------------

@app.route("/math_tutor", methods=["POST"])
@app.route("/math_tutor.php", methods=["POST"])
@limiter.limit("30 per minute")
def math_tutor():
    """
    Unified Math Tutor endpoint.

    Body:
      { "session_id": "...", "message": "...", "action": "start"|"answer"|"solve" }

    action="start"  → generate a new problem and return first step prompt
    action="answer" → validate user's step answer
    action="solve"  → direct solve mode (full solution)

    Session math_state shape:
      {
        current_problem: str,
        steps: [],
        current_step_index: int,
        final_answer: str,
        similar_problems: [],
        completed: bool
      }

    Returns:
      {
        type: "problem"|"step_feedback"|"complete"|"solution"|"error",
        problem: str,           (when type=problem)
        step_prompt: str,       (current step question)
        step_index: int,
        total_steps: int,
        correct: bool,          (on step_feedback)
        hint: str|None,
        solution_steps: [],     (on complete or direct solve)
        final_answer: str,
        similar_problems: [],
        full_solution: str      (on type=solution)
      }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    message    = str(data.get("message", "")).strip()
    action     = str(data.get("action", "start")).strip()

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    # Only operates in Math subject sessions
    if session.get("subject", "").lower() != "math":
        return jsonify({"error": "Math tutor only available for Math subject sessions"}), 400

    topic      = session["topic"]
    difficulty = session.get("depth", "Exam")

    # ── DIRECT SOLVE MODE ─────────────────────────────────────
    if action == "solve" or (action == "start" and _is_direct_solve(message) and message):
        solution = _math_direct_solve(message or topic, difficulty)
        if not solution:
            return jsonify({"error": "Could not generate solution."}), 500
        return jsonify({"type": "solution", "full_solution": solution.strip()})

    # ── START / GENERATE PROBLEM ──────────────────────────────
    if action == "start":
        prob = _generate_math_problem(topic, difficulty)
        if not prob:
            return jsonify({"error": "Could not generate problem. Please try again."}), 500

        session["math_state"] = {
            "current_problem": prob["problem"],
            "steps": prob["steps"],
            "current_step_index": 0,
            "final_answer": prob.get("final_answer", ""),
            "similar_problems": prob.get("similar_problems", []),
            "completed": False,
        }
        ms = session["math_state"]
        return jsonify({
            "type": "problem",
            "problem": ms["current_problem"],
            "step_prompt": f"Step 1: {ms['steps'][0]}",
            "step_index": 0,
            "total_steps": len(ms["steps"]),
            "final_answer": "",
            "similar_problems": [],
        })

    # ── VALIDATE STEP ANSWER ──────────────────────────────────
    if action == "answer":
        ms = session.get("math_state")
        if not ms or ms.get("completed"):
            return jsonify({"error": "No active problem. Use action=start."}), 400

        idx      = ms["current_step_index"]
        expected = ms["steps"][idx]
        result   = _validate_math_step(message, expected, ms["current_problem"], topic)

        if not result["correct"]:
            return jsonify({
                "type": "step_feedback",
                "correct": False,
                "hint": result["hint"],
                "step_index": idx,
                "total_steps": len(ms["steps"]),
                "step_prompt": f"Step {idx + 1}: Try again — {result['hint'] or 'Review your work.'}",
            })

        # Correct — advance
        idx += 1
        ms["current_step_index"] = idx

        if idx >= len(ms["steps"]):
            # Completed
            ms["completed"] = True
            return jsonify({
                "type": "complete",
                "correct": True,
                "problem": ms["current_problem"],
                "solution_steps": ms["steps"],
                "final_answer": ms["final_answer"],
                "similar_problems": ms["similar_problems"],
                "step_index": idx,
                "total_steps": len(ms["steps"]),
            })

        return jsonify({
            "type": "step_feedback",
            "correct": True,
            "hint": None,
            "step_index": idx,
            "total_steps": len(ms["steps"]),
            "step_prompt": f"Step {idx + 1}: {ms['steps'][idx]}",
        })

    return jsonify({"error": f"Unknown action: {action}"}), 400


# ---------------------------------------------------------------------------
# NEW: End-of-session topic summary
# ---------------------------------------------------------------------------

@app.route("/summary", methods=["POST"])
@app.route("/summary.php", methods=["POST"])
@limiter.limit("5 per minute")
def summary():
    """
    Generate a personalised end-of-session summary.
    Body: { "session_id": "..." }
    Returns: { "summary": "..." }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    text = generate_topic_summary(session["topic"], session["slides"], session["performance"])
    if not text:
        return jsonify({"error": "Failed to generate summary. Please try again."}), 500

    return jsonify({"summary": text.strip()})


# ---------------------------------------------------------------------------
# NEW: Per-slide learner notes
# ---------------------------------------------------------------------------

@app.route("/note", methods=["POST"])
@app.route("/note.php", methods=["POST"])
@limiter.limit("30 per minute")
def note():
    """
    Save or retrieve a learner's personal note for a slide.
    To SAVE:  { "session_id": "...", "slide_index": 0, "text": "My note..." }
    To GET:   { "session_id": "...", "slide_index": 0 }
    Returns: { "note": "..." }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    slide_index = data.get("slide_index", 0)

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    slides = session["slides"]
    if not (0 <= slide_index < len(slides)):
        return jsonify({"error": "Invalid slide index"}), 400

    cache_key = str(slide_index)

    if "text" in data:
        text = str(data["text"]).strip()[:2000]   # cap at 2000 chars
        session["notes"][cache_key] = text
        return jsonify({"note": text, "saved": True})
    else:
        return jsonify({"note": session["notes"].get(cache_key, "")})


# ---------------------------------------------------------------------------
# NEW: Per-slide sources / references
# ---------------------------------------------------------------------------

@app.route("/sources", methods=["POST"])
@app.route("/sources.php", methods=["POST"])
@limiter.limit("15 per minute")
def sources():
    """
    Generate credible reference sources for a specific slide (cached per slide).
    Body:  { "session_id": "...", "slide_index": 0 }
    Returns: { "sources": [ { title, type, authors, year, description, url_hint } ] }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    slide_index = data.get("slide_index", 0)

    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    slides = session["slides"]
    if not (0 <= slide_index < len(slides)):
        return jsonify({"error": "Invalid slide index"}), 400

    cache_key = str(slide_index)
    if cache_key in session.get("sources_cache", {}):
        return jsonify({"sources": session["sources_cache"][cache_key], "cached": True})

    result = generate_sources_for_slide(slides[slide_index], session["topic"])
    if not result:
        return jsonify({"error": "The AI could not generate sources for this slide. Please try again."}), 500

    if "sources_cache" not in session:
        session["sources_cache"] = {}
    session["sources_cache"][cache_key] = result
    return jsonify({"sources": result})


# ---------------------------------------------------------------------------
# PPT Generation
# ---------------------------------------------------------------------------

def build_pptx(topic: str, slides_data: list, theme: str = "dark"):
    """
    Build a beautifully designed PPTX in either dark or light theme.
    Returns a python-pptx Presentation object, or None if pptx is unavailable.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree
    except ImportError:
        return None

    is_dark = (theme != "light")

    # ── Colour palettes ─────────────────────────────────────────
    if is_dark:
        C_BG      = RGBColor(0x08, 0x0b, 0x12)   # near-black
        C_SURF    = RGBColor(0x0f, 0x14, 0x20)   # surface
        C_CARD    = RGBColor(0x14, 0x19, 0x26)   # card
        C_BORDER  = RGBColor(0x2a, 0x33, 0x47)   # border
        C_BORDER2 = RGBColor(0x1e, 0x25, 0x35)
        C_ACCENT1 = RGBColor(0xe8, 0xa0, 0x20)   # gold
        C_ACCENT2 = RGBColor(0x20, 0xc5, 0xb0)   # teal
        C_TITLE   = RGBColor(0xea, 0xe6, 0xdf)   # near-white
        C_BODY    = RGBColor(0xb8, 0xc0, 0xd4)   # light grey body text
        C_MUTED   = RGBColor(0x8b, 0x93, 0xa8)
        C_DIM     = RGBColor(0x4a, 0x52, 0x68)
        C_PILL_BG = RGBColor(0x0b, 0x1f, 0x1c)
        BULLET_COLORS = [C_ACCENT1, C_ACCENT2, C_ACCENT1, C_ACCENT2, C_ACCENT1]
    else:
        C_BG      = RGBColor(0xf7, 0xf8, 0xfc)   # off-white
        C_SURF    = RGBColor(0xff, 0xff, 0xff)   # white
        C_CARD    = RGBColor(0xee, 0xf0, 0xf8)   # light card
        C_BORDER  = RGBColor(0xd0, 0xd6, 0xe8)
        C_BORDER2 = RGBColor(0xe0, 0xe5, 0xf2)
        C_ACCENT1 = RGBColor(0xc4, 0x7d, 0x0e)   # amber (darker gold for contrast)
        C_ACCENT2 = RGBColor(0x0d, 0x9c, 0x8a)   # dark teal
        C_TITLE   = RGBColor(0x12, 0x17, 0x2e)   # near-black
        C_BODY    = RGBColor(0x2c, 0x35, 0x52)   # dark navy body
        C_MUTED   = RGBColor(0x60, 0x6a, 0x88)
        C_DIM     = RGBColor(0x90, 0x9a, 0xb8)
        C_PILL_BG = RGBColor(0xe0, 0xf5, 0xf2)
        BULLET_COLORS = [C_ACCENT1, C_ACCENT2, C_ACCENT1, C_ACCENT2, C_ACCENT1]

    W = Inches(13.333)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ── Helpers ─────────────────────────────────────────────────

    def set_bg(slide, color):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def rect(slide, l, t, w, h, fill, line=None, line_w=Pt(0), rounding=None):
        shp = slide.shapes.add_shape(1, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
            shp.line.width = line_w
        else:
            shp.line.fill.background()
        return shp

    def textbox(slide, text, l, t, w, h,
                font='Calibri', size=12, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name  = font
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return tb

    # Expand long slides into continuation parts so full sentences fit in PPT.
    expanded_slides = []
    for s in slides_data:
        pts = s.get("points", []) if isinstance(s, dict) else []
        if not isinstance(pts, list):
            pts = []
        point_texts = [
            str(p.get("text", "") if isinstance(p, dict) else p).strip()
            for p in pts
            if str(p.get("text", "") if isinstance(p, dict) else p).strip()
        ]
        max_len = max((len(t) for t in point_texts), default=0)
        if max_len > 420:
            chunk_size = 2
        elif max_len > 260:
            chunk_size = 3
        elif max_len > 170:
            chunk_size = 4
        else:
            chunk_size = 5
        if not pts:
            expanded_slides.append({**s, "_ppt_points": [], "_ppt_part": 1, "_ppt_parts": 1})
            continue
        total_parts = (len(pts) + chunk_size - 1) // chunk_size
        for i in range(0, len(pts), chunk_size):
            expanded_slides.append({
                **s,
                "_ppt_points": pts[i:i + chunk_size],
                "_ppt_part": (i // chunk_size) + 1,
                "_ppt_parts": total_parts,
            })

    total = len(expanded_slides)

    # ════════════════════════════════════════════════════════════
    # SLIDE 0 — TITLE SLIDE
    # ════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    set_bg(sl, C_BG)

    # ── Full-height left panel ──────────────────────────────────
    PANEL_W = Inches(5.2)
    rect(sl, 0, 0, PANEL_W, H, C_SURF)
    rect(sl, PANEL_W, 0, Inches(0.032), H, C_BORDER)

    # Top accent bar inside left panel
    rect(sl, 0, 0, PANEL_W, Inches(0.22), C_ACCENT1)

    # Decorative stacked rectangles (visual layering)
    rect(sl, Inches(0.35), Inches(0.55), Inches(0.9), Inches(0.9), C_CARD,
         line=C_BORDER, line_w=Pt(0.5))
    rect(sl, Inches(0.55), Inches(0.75), Inches(0.9), Inches(0.9), C_ACCENT1)
    textbox(sl, 'AI', Inches(0.6), Inches(0.78), Inches(0.8), Inches(0.85),
            font='Georgia', size=20, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)

    # Brand label
    textbox(sl, 'ADAPTIVE LEARN  ·  AI LECTURE ENGINE',
            Inches(0.35), Inches(1.85), PANEL_W - Inches(0.7), Inches(0.3),
            font='Courier New', size=7.5, bold=True, color=C_ACCENT2)

    # Thin rule
    rect(sl, Inches(0.35), Inches(2.22), Inches(0.55), Inches(0.038), C_ACCENT1)

    # Main title text
    t_short = topic if len(topic) <= 55 else topic[:53] + '…'
    tb = sl.shapes.add_textbox(Inches(0.35), Inches(2.35), PANEL_W - Inches(0.7), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = t_short
    run.font.name  = 'Georgia'
    run.font.size  = Pt(34)
    run.font.bold  = True
    run.font.color.rgb = C_TITLE

    # Slide count badge
    rect(sl, Inches(0.35), Inches(5.5), Inches(1.6), Inches(0.5), C_ACCENT1)
    textbox(sl, f'{total}  SLIDES', Inches(0.35), Inches(5.52),
            Inches(1.6), Inches(0.46),
            font='Courier New', size=11, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)

    textbox(sl, f'Generated by AdaptiveLearn AI',
            Inches(0.35), Inches(6.15), PANEL_W - Inches(0.7), Inches(0.3),
            font='Calibri', size=10, italic=True, color=C_MUTED)

    # ── Right decorative panel ──────────────────────────────────
    RX = PANEL_W + Inches(0.032)
    RW = W - RX

    # Large accent block top-right
    rect(sl, RX, 0, RW, Inches(2.5), C_ACCENT1)
    # Teal strip below it
    rect(sl, RX, Inches(2.5), RW, Inches(0.18), C_ACCENT2)

    # Big decorative number
    textbox(sl, str(total), RX, Inches(0.05), RW, Inches(2.4),
            font='Georgia', size=110, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)

    # Grid of decorative cards in right panel body
    card_y_start = Inches(2.9)
    card_h = Inches(1.1)
    card_gap = Inches(0.18)
    deco_labels = ['CONCEPTS', 'QUIZZES', 'INSIGHTS', 'REVIEW']
    for ci in range(min(4, total)):
        cy = card_y_start + ci * (card_h + card_gap)
        rect(sl, RX + Inches(0.3), cy, RW - Inches(0.6), card_h,
             C_CARD if is_dark else C_SURF,
             line=C_BORDER, line_w=Pt(0.6))
        rect(sl, RX + Inches(0.3), cy, Inches(0.06), card_h, C_ACCENT2)
        label = deco_labels[ci] if ci < len(deco_labels) else f'PART {ci+1}'
        textbox(sl, label, RX + Inches(0.5), cy + Inches(0.15),
                RW - Inches(1.0), Inches(0.3),
                font='Courier New', size=7, bold=True, color=C_ACCENT2)
        textbox(sl, f'Section {ci + 1}', RX + Inches(0.5), cy + Inches(0.45),
                RW - Inches(1.0), Inches(0.4),
                font='Calibri', size=12, bold=True, color=C_TITLE)

    # Bottom strip
    rect(sl, 0, H - Inches(0.12), W, Inches(0.12), C_ACCENT2)

    # ════════════════════════════════════════════════════════════
    # CONTENT SLIDES
    # ════════════════════════════════════════════════════════════
    for idx, sdata in enumerate(expanded_slides):
        sl = prs.slides.add_slide(blank)
        set_bg(sl, C_BG)

        title_txt = str(sdata.get('title', f'Slide {idx + 1}')).strip()
        part = int(sdata.get("_ppt_part", 1))
        parts = int(sdata.get("_ppt_parts", 1))
        if parts > 1:
            title_txt = f"{title_txt} (Part {part}/{parts})"
        points_raw = sdata.get('_ppt_points', sdata.get('points', []))
        n_pts = min(len(points_raw), 5)

        # ── Top header band ───────────────────────────────────────
        HEADER_H = Inches(1.45)
        rect(sl, 0, 0, W, HEADER_H, C_SURF)

        # Gold top stripe
        rect(sl, 0, 0, W, Inches(0.18), C_ACCENT1)

        # Teal left stripe (full height)
        rect(sl, 0, 0, Inches(0.1), H, C_ACCENT2)

        # Slide number box — top-left
        rect(sl, Inches(0.22), Inches(0.35), Inches(1.15), Inches(0.72), C_ACCENT1)
        textbox(sl, f'{idx + 1:02d}', Inches(0.22), Inches(0.35),
                Inches(1.15), Inches(0.72),
                font='Georgia', size=26, bold=True, color=C_BG,
                align=PP_ALIGN.CENTER)

        # Slide counter label
        textbox(sl, f'/ {total}', Inches(1.42), Inches(0.62),
                Inches(0.7), Inches(0.32),
                font='Courier New', size=9, color=C_MUTED)

        # Topic label right
        t_label = (topic[:52] + '…') if len(topic) > 52 else topic
        textbox(sl, t_label.upper(),
                W - Inches(5.5), Inches(0.55), Inches(5.2), Inches(0.28),
                font='Courier New', size=7.5, color=C_DIM,
                align=PP_ALIGN.RIGHT)

        # Separator line below header
        rect(sl, Inches(0.1), HEADER_H, W - Inches(0.1), Inches(0.025), C_BORDER)

        # ── Slide title ──────────────────────────────────────────
        title_disp = title_txt
        tb = sl.shapes.add_textbox(Inches(0.3), Inches(1.55), W - Inches(2.2), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_disp
        run.font.name  = 'Georgia'
        run.font.size  = Pt(26)
        run.font.bold  = True
        run.font.color.rgb = C_TITLE

        # Underline accent beneath title
        rect(sl, Inches(0.3), Inches(2.54), Inches(0.6), Inches(0.04), C_ACCENT1)
        rect(sl, Inches(0.95), Inches(2.54), Inches(0.22), Inches(0.04), C_ACCENT2)

        # ── Right sidebar ─────────────────────────────────────────
        SB_X = W - Inches(1.9)
        SB_W = Inches(1.75)
        rect(sl, SB_X, HEADER_H + Inches(0.08), SB_W, H - HEADER_H - Inches(0.35),
             C_CARD if is_dark else C_SURF,
             line=C_BORDER, line_w=Pt(0.5))
        rect(sl, SB_X, HEADER_H + Inches(0.08), Inches(0.06),
             H - HEADER_H - Inches(0.35), C_ACCENT2)

        # Sidebar: big slide number watermark
        textbox(sl, str(idx + 1),
                SB_X, Inches(3.0), SB_W, Inches(1.4),
                font='Georgia', size=64, bold=True, color=C_BORDER,
                align=PP_ALIGN.CENTER)

        # Sidebar: section label
        textbox(sl, 'SLIDE', SB_X, Inches(2.4), SB_W, Inches(0.35),
                font='Courier New', size=8, bold=True, color=C_DIM,
                align=PP_ALIGN.CENTER)

        # Sidebar: topic dots (decorative)
        for di in range(min(n_pts, 5)):
            dot_y = Inches(4.6) + di * Inches(0.38)
            rect(sl, SB_X + Inches(0.75), dot_y, Inches(0.12), Inches(0.12),
                 C_ACCENT1 if di % 2 == 0 else C_ACCENT2)

        # ── Bullet points area ────────────────────────────────────
        font_size = 13 if n_pts <= 3 else (12 if n_pts == 4 else 11)
        bullet_h  = Inches(0.78) if n_pts <= 3 else (Inches(0.7) if n_pts == 4 else Inches(0.62))
        start_y   = Inches(2.72)
        content_w = SB_X - Inches(0.45)

        for bi, point in enumerate(points_raw[:5]):
            pt_text = str(point.get('text', '') if isinstance(point, dict) else point).strip()
            if not pt_text:
                continue

            by = start_y + bi * bullet_h
            bc = BULLET_COLORS[bi % len(BULLET_COLORS)]

            # Card-style row background (alternating subtle tint)
            row_bg = (RGBColor(0x10, 0x16, 0x24) if is_dark else RGBColor(0xf0, 0xf3, 0xfc)) \
                     if bi % 2 == 0 else \
                     (RGBColor(0x0d, 0x12, 0x1e) if is_dark else RGBColor(0xfa, 0xfb, 0xff))
            rect(sl, Inches(0.22), by, content_w - Inches(0.22), bullet_h - Inches(0.06),
                 row_bg, line=C_BORDER2, line_w=Pt(0.4))

            # Left colour stripe on each row
            rect(sl, Inches(0.22), by, Inches(0.06), bullet_h - Inches(0.06), bc)

            # Bullet number circle (simulated with text)
            rect(sl, Inches(0.38), by + Inches(0.19), Inches(0.32), Inches(0.32), bc)
            textbox(sl, str(bi + 1),
                    Inches(0.38), by + Inches(0.18), Inches(0.32), Inches(0.34),
                    font='Courier New', size=9, bold=True, color=C_BG,
                    align=PP_ALIGN.CENTER)

            # Point text (full sentence; no truncation)
            pt_disp = pt_text
            tb2 = sl.shapes.add_textbox(Inches(0.82), by + Inches(0.08),
                                         content_w - Inches(1.1), bullet_h - Inches(0.14))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = pt_disp
            run2.font.name  = 'Calibri'
            text_len = len(pt_disp)
            adj_size = font_size
            if text_len > 420:
                adj_size = max(8, font_size - 3)
            elif text_len > 260:
                adj_size = max(9, font_size - 2)
            elif text_len > 170:
                adj_size = max(10, font_size - 1)
            run2.font.size = Pt(adj_size)
            run2.font.color.rgb = C_BODY

        # ── Bottom bar ────────────────────────────────────────────
        rect(sl, 0, H - Inches(0.14), W, Inches(0.14), C_ACCENT1)
        textbox(sl, 'ADAPTIVELEARN  ·  AI-POWERED EDUCATION',
                Inches(0.22), H - Inches(0.13), Inches(5.0), Inches(0.13),
                font='Courier New', size=6.5, bold=True, color=C_BG)
        textbox(sl, f'SLIDE {idx + 1} OF {total}',
                W - Inches(2.2), H - Inches(0.13), Inches(2.0), Inches(0.13),
                font='Courier New', size=6.5, bold=True, color=C_BG,
                align=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════════════════════
    # CLOSING SLIDE
    # ════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    set_bg(sl, C_BG)

    # Bold split — left teal panel, right body
    rect(sl, 0, 0, Inches(4.8), H, C_ACCENT2)
    rect(sl, Inches(4.8), 0, Inches(0.06), H, C_ACCENT1)

    # Large checkmark-ish decoration on left panel
    textbox(sl, '✓', Inches(0.3), Inches(1.1), Inches(4.2), Inches(2.0),
            font='Georgia', size=100, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)
    textbox(sl, 'SESSION COMPLETE', Inches(0.3), Inches(3.3), Inches(4.2), Inches(0.55),
            font='Courier New', size=13, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)
    textbox(sl, f'{total} slides covered', Inches(0.3), Inches(3.95),
            Inches(4.2), Inches(0.4),
            font='Calibri', size=11, italic=True, color=RGBColor(0xd0,0xf5,0xef),
            align=PP_ALIGN.CENTER)

    # Right panel
    t_close = topic if len(topic) <= 58 else topic[:56] + '…'
    textbox(sl, 'YOU HAVE STUDIED', Inches(5.3), Inches(1.5),
            W - Inches(5.7), Inches(0.4),
            font='Courier New', size=9, bold=True, color=C_ACCENT2)
    rect(sl, Inches(5.3), Inches(1.97), Inches(0.55), Inches(0.038), C_ACCENT1)

    tb = sl.shapes.add_textbox(Inches(5.3), Inches(2.1), W - Inches(5.7), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = t_close
    run.font.name  = 'Georgia'
    run.font.size  = Pt(36)
    run.font.bold  = True
    run.font.color.rgb = C_TITLE

    textbox(sl, f'Powered by AdaptiveLearn AI  ·  {total} slides',
            Inches(5.3), Inches(5.0), W - Inches(5.7), Inches(0.4),
            font='Calibri', size=11, italic=True, color=C_MUTED)

    # Bottom bar
    rect(sl, 0, H - Inches(0.14), W, Inches(0.14), C_ACCENT1)

    return prs


@app.route("/generate_ppt", methods=["POST"])
@app.route("/generate_ppt.php", methods=["POST"])
@limiter.limit("5 per minute")
def generate_ppt():
    """
    Build and return a .pptx file for the current session.
    Body: { "session_id": "...", "theme": "dark"|"light" }
    Returns: binary .pptx file download
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    session_id = data.get("session_id", "")
    theme      = data.get("theme", "dark")
    session = sessions.get(session_id)
    if not session:
        return jsonify({"error": "Session not found"}), 404

    prs = build_pptx(session["topic"], session["slides"], theme=theme)
    if prs is None:
        return jsonify({
            "error": "python-pptx is not installed on the server. "
                     "Run: pip install python-pptx"
        }), 500

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    safe_name = re.sub(r'[^a-zA-Z0-9_-]', '_', session["topic"][:40]).strip('_')
    theme_tag  = "Dark" if theme != "light" else "Light"
    filename = f"{safe_name}_{theme_tag}_AdaptiveLearn.pptx"

    return send_file(
        buf,
        as_attachment=True,
        download_name=filename,
        mimetype=(
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
    )


# ---------------------------------------------------------------------------
# _ASSIGNMENT — Law Assignment Generator (inlined from blueprint)
# ---------------------------------------------------------------------------
# Token budgets: each section ~800-1200 tokens to guarantee ≥1 full page of
# content per section. 6 sections × ~2.5 pages ≈ 15 pages total.
# Every section is a SEPARATE _call_groq call so output never corrupts midway.
# ---------------------------------------------------------------------------

import io as _io
from datetime import datetime as _dt

_ASSIGN_STORE: dict = {}      # { id: { meta, sections, created_at } }
_ASSIGN_TTL   = 7200          # 2 hours

# Token budget per section — tuned for ≥1 full A4 page each (~600-800 words)
_ASSIGN_TOKEN = {
    "declaration":              350,    # ~0.5 page
    "acknowledgment":           400,    # ~0.6 page
    "introduction":            1200,    # ~2 pages
    "concept":                 1100,    # ~1.5 pages — core doctrinal concept
    "structural_functional":    900,    # ~1.2 pages
    "checks_balances":         1000,    # ~1.5 pages
    "india_context":            900,    # ~1.2 pages
    "assembly_debates":         900,    # ~1.2 pages
    "constitutional_reflection":900,    # ~1.2 pages
    "case_law":                1200,    # ~1.8 pages — landmark cases
    "conclusion":               800,    # ~1.2 pages
    "bibliography":             700,    # ~1 page
}


def _assign_prune():
    now = time.time()
    stale = [k for k, v in _ASSIGN_STORE.items() if now - v["created_at"] > _ASSIGN_TTL]
    for k in stale:
        del _ASSIGN_STORE[k]


def _assign_find_cached(topic: str):
    for v in _ASSIGN_STORE.values():
        if v["meta"]["topic"].lower().strip() == topic.lower().strip():
            return v
    return None


def _assign_prompt(section: str, topic: str, meta: dict) -> str:
    """Build a detailed prompt per section. Each gets its own API call.
    Supports 13-section CUSB-style academic assignments targeting ~15 pages of content."""

    # Build author/book instruction if provided
    ref_authors = meta.get("reference_authors", [])
    if ref_authors:
        author_lines = []
        for e in ref_authors:
            line = e["author"]
            if e.get("book"):
                book_title = e["book"]
                line += f", \"{book_title}\""
            author_lines.append(line)
        author_instruction = (
            "IMPORTANT: Draw content, arguments, and analysis primarily from the following "
            "authors and their works (cite them specifically in the text where relevant): "
            + "; ".join(author_lines) + ". "
        )
    else:
        author_instruction = ""

    if section == "declaration":
        return (
            f"Write a formal academic honesty declaration for a law assignment. "
            f"Student: {meta['student_name']}, Roll No: {meta['roll_no']}, "
            f"Course: {meta['course']}, College: {meta['college']}. "
            f"The declaration must be at least 180 words. Include: "
            f"(1) Statement of original authorship, (2) Confirmation no plagiarism, "
            f"(3) Sources properly acknowledged, (4) No unauthorized assistance, "
            f"(5) Understanding of academic consequences. Formal, first-person tone. "
            f"End with student name and college name on separate lines."
        )

    elif section == "acknowledgment":
        return (
            f"Write a detailed acknowledgment (minimum 200 words) for a law assignment by "
            f"{meta['student_name']} of {meta['college']} on the topic: \"{topic}\". "
            f"Acknowledge: (1) H.O.D. and department head with formal gratitude, "
            f"(2) Subject teacher / faculty supervisor for specific guidance on this topic, "
            f"(3) College library and research resources, "
            f"(4) Family support and encouragement, "
            f"(5) Friends and classmates for peer discussions. "
            f"Warm but formal academic tone. Close with 'Thanking You' and student name."
        )

    elif section == "introduction":
        return (
            f"{author_instruction}Write a comprehensive Introduction (minimum 600 words) for a law assignment on: \"{topic}\". "
            f"Structure with these elements: "
            f"(1) Opening hook — quote a relevant jurist, philosopher, or constitutional provision to set the context; "
            f"(2) Clear definition of key legal terms and the core doctrine involved; "
            f"(3) Historical origins — trace the doctrine from ancient thinkers (e.g. Aristotle, Locke, Montesquieu) "
            f"to its modern form, naming specific works and dates; "
            f"(4) Constitutional or statutory basis in Indian law — cite relevant Articles; "
            f"(5) Scope and limitations of the present study; "
            f"(6) Significance and relevance in contemporary legal context; "
            f"(7) Thesis statement summarising the assignment's central argument. "
            f"Use flowing paragraphs. Never use ** or ## symbols. Formal academic legal prose."
        )

    elif section == "concept":
        return (
            f"{author_instruction}Write a detailed section titled 'THE CONCEPT AND THEORETICAL FOUNDATIONS' "
            f"(minimum 550 words) for a law assignment on: \"{topic}\". "
            f"Cover the following comprehensively: "
            f"(1) Core theoretical meaning and doctrinal definition of the topic with reference to major jurists; "
            f"(2) Philosophical justification — why this doctrine is necessary in a constitutional democracy, "
            f"drawing on thinkers such as Locke, Montesquieu, Madison, or other relevant philosophers; "
            f"(3) Different scholarly interpretations and academic debates around the concept; "
            f"(4) Distinction between the strict/pure form of the doctrine and its diluted/practical application; "
            f"(5) How the doctrine functions as a safeguard against tyranny and arbitrariness. "
            f"Write in full formal paragraphs. Cite scholars and jurists by name where possible."
        )

    elif section == "structural_functional":
        return (
            f"{author_instruction}Write a detailed section titled 'STRUCTURAL VERSUS FUNCTIONAL DIMENSION' "
            f"(minimum 500 words) for a law assignment on: \"{topic}\". "
            f"Analyse: "
            f"(1) What it means for the doctrine to be 'structural rather than functional' — "
            f"the idea that separation is about institutional design, not just task allocation; "
            f"(2) The three principal organs of government (Legislature, Executive, Judiciary) — "
            f"their distinct structural roles and how they interact under this doctrine; "
            f"(3) The problem of functional overlap in modern governance — delegated legislation, "
            f"quasi-judicial executive bodies, tribunals, ordinance-making powers; "
            f"(4) Academic critiques of rigid structural separation and arguments for flexible functional separation; "
            f"(5) How the Indian constitutional design reflects a structural but not absolute separation. "
            f"Use Lord Acton's maxim on power where relevant. Formal academic tone."
        )

    elif section == "checks_balances":
        return (
            f"{author_instruction}Write a detailed section titled 'THEORY OF CHECKS AND BALANCES' "
            f"(minimum 550 words) for a law assignment on: \"{topic}\". "
            f"Cover: "
            f"(1) Historical origin of the checks and balances theory — Polybius, Roman mixed government, "
            f"17th-century political theorists; "
            f"(2) The Madisonian Model — how the U.S. framers operationalised checks and balances "
            f"alongside separation of powers (presidential veto, Senate confirmation, judicial review); "
            f"(3) The two foundational principles: (a) power corrupts — no organ should have unchecked power; "
            f"(b) power must be checked by power — mutual restraint between branches; "
            f"(4) How checks and balances supplement separation of powers — the inter-organ relationship; "
            f"(5) Application in the Indian context — legislative override, executive appointment of judges, "
            f"judicial review of legislation, impeachment. "
            f"Specific examples from Indian constitutional practice required. Formal prose throughout."
        )

    elif section == "india_context":
        return (
            f"{author_instruction}Write a detailed section titled 'SEPARATION OF POWERS IN INDIA — THE CONSTITUTIONAL POSITION' "
            f"(minimum 500 words) for a law assignment on: \"{topic}\". "
            f"Analyse: "
            f"(1) Whether the Indian Constitution expressly adopts the doctrine — explain why there is no rigid separation "
            f"but sufficient functional differentiation (cite Articles 50, 121, 122, 211, 212, 361 as relevant to the topic); "
            f"(2) Executive power: President and Governor under Articles 52-78, 154, 161; "
            f"(3) Legislative power: Parliament's supremacy under Articles 107-122, and limits thereon; "
            f"(4) Judicial independence: Articles 124-147 and protection of judiciary from executive interference; "
            f"(5) Overlapping functions that exist by constitutional design — e.g. President's ordinance power (Art. 123), "
            f"Money Bills, Parliamentary privileges; "
            f"(6) How the Indian parliamentary system differs from the American presidential model of separation. "
            f"Cite specific Articles throughout. Formal academic tone."
        )

    elif section == "assembly_debates":
        return (
            f"{author_instruction}Write a detailed section titled 'THE CONSTITUENT ASSEMBLY DEBATES' "
            f"(minimum 500 words) for a law assignment on: \"{topic}\". "
            f"Reconstruct the key debates in the Constituent Assembly about how the doctrine of separation of powers "
            f"and related principles were considered during constitution-drafting. Cover: "
            f"(1) The proposal to include an express separation of powers clause — who proposed it and the reasoning; "
            f"(2) Arguments in favour of strict separation (presidential model camp); "
            f"(3) Arguments against rigid separation — Dr. B.R. Ambedkar's position favouring parliamentary system "
            f"with cooperation between executive and legislature; "
            f"(4) The final decision and why the Constituent Assembly rejected an express separation of powers clause; "
            f"(5) What this debate reveals about the framers' vision for the Indian constitution as it relates to: \"{topic}\". "
            f"Write in formal analytical prose. Use specific names of Assembly members where possible."
        )

    elif section == "constitutional_reflection":
        return (
            f"{author_instruction}Write a detailed section titled 'REFLECTION OF THE DOCTRINE IN THE INDIAN CONSTITUTION' "
            f"(minimum 500 words) for a law assignment on: \"{topic}\". "
            f"Demonstrate how, even without express provision, the doctrine is reflected in the Constitution: "
            f"(1) Distribution of powers between Legislature, Executive, Judiciary — each having exclusive core functions; "
            f"(2) Fixed tenure, enumerated powers, and constitutional salaries protecting institutional independence; "
            f"(3) Judicial review under Articles 13, 32, 226 as the primary check on legislative and executive action; "
            f"(4) Parliamentary controls on the executive through question hour, no-confidence motions, budget approval; "
            f"(5) Constitutional bars preventing discussion of judicial conduct in Parliament (Art. 121); "
            f"(6) The basic structure doctrine as a supreme-level constraint on Parliament's amending power — "
            f"how this represents the highest form of separation of judicial power from legislative supremacy. "
            f"Formal academic prose. Cite relevant Articles by number."
        )

    elif section == "case_law":
        return (
            f"{author_instruction}Write a detailed section titled 'LANDMARK CASE LAW AND JUDICIAL INTERPRETATION' "
            f"(minimum 700 words) for a law assignment on: \"{topic}\". "
            f"Analyse at least 6 landmark Indian Supreme Court judgments relevant to: \"{topic}\". "
            f"For each case provide: (a) full case name and year, (b) brief facts, (c) key legal question, "
            f"(d) court's holding and the specific observation about separation of powers or the relevant doctrine, "
            f"(e) significance and impact on subsequent law. "
            f"Mandatory cases to cover (adapt if topic differs): "
            f"(1) Ram Jawaya Kapur v. State of Punjab (1955) — on absence of strict separation in India; "
            f"(2) Kesavananda Bharati v. State of Kerala (1973) — basic structure and judicial power; "
            f"(3) Indira Nehru Gandhi v. Raj Narain (1975) — separation in electoral disputes; "
            f"(4) Minerva Mills v. Union of India (1980) — limits on Parliament's amending power; "
            f"(5) Delhi Laws Act case — delegation of legislative power; "
            f"(6) Any recent relevant Supreme Court judgment post-2010. "
            f"If the topic is not separation of powers, substitute with the most relevant landmark cases for: \"{topic}\". "
            f"Formal analytical prose. No bullet points — write in full paragraphs."
        )

    elif section == "conclusion":
        return (
            f"Write a comprehensive Conclusion (minimum 450 words) for a law assignment on: \"{topic}\". "
            f"Include: "
            f"(1) Synthesis of key findings from each section — do not introduce new arguments; "
            f"(2) Restatement of the thesis in light of the analysis; "
            f"(3) The broader constitutional significance of the topic for Indian democracy and rule of law; "
            f"(4) At least 3 specific, actionable policy recommendations or suggestions for legal reform; "
            f"(5) Identification of areas requiring further judicial or legislative attention; "
            f"(6) Closing reflection on why the topic remains vital in contemporary constitutional governance. "
            f"Formal, measured academic tone. No markdown. No bullet points in the actual text — flowing paragraphs."
        )

    elif section == "bibliography":
        return (
            f"{author_instruction}Write a complete Bibliography in Bluebook citation format for a law assignment on: \"{topic}\". "
            f"If reference authors were specified above, ensure they appear as primary entries in the Books section. "
            f"Provide at least 12 entries organised in these categories: "
            f"A. Cases (minimum 4 entries — full citation with year and court); "
            f"B. Statutes and Constitutional Provisions (minimum 2 entries — e.g. Constitution of India Articles, relevant Acts); "
            f"C. Books and Textbooks (minimum 3 entries — author, title in italics, edition, publisher, year); "
            f"D. Journal Articles (minimum 2 entries — author, title, volume, journal name, year, page numbers); "
            f"E. Online Resources (minimum 1 entry — author if available, title, URL, last visited date). "
            f"Number all entries within each category. "
            f"Format strictly in Bluebook style. Mark any uncertain citations as [CITATION NEEDED]."
        )

    return f"Write a detailed section on {section} for a law assignment about: \"{topic}\"."


def _assign_gen_section(section: str, topic: str, meta: dict) -> str:
    """Each section = one dedicated _call_groq call. No batching = no corruption."""
    system = (
        "You are a senior legal academic writer at a reputed Indian law school. "
        "Write in formal, precise academic English. Use correct legal terminology. "
        "Structure content clearly. Never use markdown symbols like ** or ##. "
        "Write in full sentences and flowing paragraphs. Minimum length per your instructions."
    )
    result = _call_groq(
        _assign_prompt(section, topic, meta),
        max_tokens=_ASSIGN_TOKEN.get(section, 800),
        system=system,
    )
    return (result or "").strip()


def _assign_assemble(meta: dict) -> dict:
    _assign_prune()
    cached = _assign_find_cached(meta["topic"])
    if cached:
        return cached

    sections: dict = {}
    order = ["declaration"]
    if meta.get("include_acknowledgment", True):
        order.append("acknowledgment")
    order += [
        "introduction",
        "concept",
        "structural_functional",
        "checks_balances",
        "india_context",
        "assembly_debates",
        "constitutional_reflection",
        "case_law",
        "conclusion",
        "bibliography",
    ]

    # SEPARATE API call per section — prevents mid-response corruption
    for sec in order:
        sections[sec] = _assign_gen_section(sec, meta["topic"], meta)

    record = {
        "id":         str(uuid.uuid4()),
        "meta":       meta,
        "sections":   sections,
        "created_at": time.time(),
    }
    _ASSIGN_STORE[record["id"]] = record
    return record


# ── DOCX builder ─────────────────────────────────────────────────────────────

def _assign_build_docx(record: dict) -> bytes:
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")

    meta     = record["meta"]
    sections = record["sections"]
    doc      = Document()

    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.25)
        section.right_margin  = Inches(1.0)

    RED = RGBColor(0xCC, 0x00, 0x00)

    def set_font(run, size=12, bold=False, italic=False, color=None, underline=False):
        run.font.name    = "Times New Roman"
        run.font.size    = Pt(size)
        run.bold         = bold
        run.italic       = italic
        run.underline    = underline
        if color:
            run.font.color.rgb = color

    def add_heading(text, level=1):
        """Bold + underlined + red heading for every section title."""
        p   = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(18 if level == 1 else 12)
        p.paragraph_format.space_after  = Pt(8)
        run = p.add_run(text)
        set_font(run, size=14 if level == 1 else 12,
                 bold=True, underline=True, color=RED)
        p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
        return p

    def add_subheading(text):
        """Sub-section headings (2.1, 2.2 …) — bold + underlined + red, smaller."""
        p   = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after  = Pt(5)
        run = p.add_run(text)
        set_font(run, size=12, bold=True, underline=True, color=RED)
        p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
        return p

    def add_body(text):
        if not text:
            return
        for line in text.split("\n"):
            line = line.strip()
            if not line:
                continue
            is_subhead = bool(re.match(r'^(\d+\.\d+|[IVX]+\.)\s', line))
            if is_subhead:
                add_subheading(line)
            else:
                p   = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(6)
                run = p.add_run(line)
                set_font(run, size=12)
                p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    def add_divider():
        p   = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after  = Pt(4)
        pPr    = p._p.get_or_add_pPr()
        pBdr   = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"),   "single")
        bottom.set(qn("w:sz"),    "6")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "CC0000")   # red divider
        pBdr.append(bottom)
        pPr.append(pBdr)

    def page_break():
        doc.add_page_break()

    # ── COVER PAGE ──────────────────────────────────────────────────────────
    for _ in range(4):
        doc.add_paragraph()
    p   = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(meta["college"].upper())
    set_font(run, 16, bold=True, color=RGBColor(0x1A, 0x3A, 0x5C))

    p   = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(meta["course"])
    set_font(run, 13)

    doc.add_paragraph()
    add_divider()
    doc.add_paragraph()

    p   = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("LAW ASSIGNMENT")
    set_font(run, 20, bold=True, color=RGBColor(0x1A, 0x3A, 0x5C))

    p   = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("ON")
    set_font(run, 12)

    p   = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(meta["topic"].upper())
    set_font(run, 15, bold=True, color=RGBColor(0x1A, 0x3A, 0x5C))

    doc.add_paragraph()
    add_divider()
    for _ in range(3):
        doc.add_paragraph()

    for label, value in [
        ("Submitted by", meta["student_name"]),
        ("Roll No",      meta["roll_no"]),
        ("Date",         _dt.now().strftime("%d %B %Y")),
    ]:
        p   = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r1  = p.add_run(f"{label}: ")
        set_font(r1, 12)
        r2  = p.add_run(value)
        set_font(r2, 12, bold=True)

    page_break()

    # ── TABLE OF CONTENTS ────────────────────────────────────────────────────
    add_heading("TABLE OF CONTENTS")
    add_divider()
    toc_entries = [
        ("Declaration",                                    "2"),
        ("Acknowledgment",                                 "3") if sections.get("acknowledgment") else None,
        ("1.  Introduction",                               "4"),
        ("2.  The Concept and Theoretical Foundations",    "6"),
        ("3.  Structural Versus Functional Dimension",     "7"),
        ("4.  Theory of Checks and Balances",              "9"),
        ("5.  Separation of Powers in India",              "11"),
        ("6.  The Constituent Assembly Debates",           "12"),
        ("7.  Reflection of the Doctrine in Constitution", "13"),
        ("8.  Landmark Case Law",                          "14"),
        ("9.  Conclusion",                                 "16"),
        ("10. Bibliography",                               "17"),
    ]
    for entry in toc_entries:
        if not entry:
            continue
        p   = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        r1  = p.add_run(entry[0])
        set_font(r1, 12)
        r2  = p.add_run(f"{'.' * max(1, 55 - len(entry[0]))} {entry[1]}")
        set_font(r2, 12)
    page_break()

    # ── DECLARATION ──────────────────────────────────────────────────────────
    add_heading("DECLARATION")
    add_divider()
    add_body(sections.get("declaration", ""))
    doc.add_paragraph()
    p   = doc.add_paragraph()
    run = p.add_run(f"Signature: ________________        Date: {_dt.now().strftime('%d %B %Y')}")
    set_font(run, 11)
    page_break()

    # ── ACKNOWLEDGMENT ───────────────────────────────────────────────────────
    if sections.get("acknowledgment"):
        add_heading("ACKNOWLEDGMENT")
        add_divider()
        add_body(sections["acknowledgment"])
        page_break()

    # ── INTRODUCTION ─────────────────────────────────────────────────────────
    add_heading("1.  INTRODUCTION")
    add_divider()
    add_body(sections.get("introduction", ""))
    page_break()

    # ── CONCEPT ──────────────────────────────────────────────────────────────
    add_heading("2.  THE CONCEPT AND THEORETICAL FOUNDATIONS")
    add_divider()
    add_body(sections.get("concept", ""))
    page_break()

    # ── STRUCTURAL vs FUNCTIONAL ──────────────────────────────────────────────
    add_heading("3.  STRUCTURAL VERSUS FUNCTIONAL DIMENSION")
    add_divider()
    add_body(sections.get("structural_functional", ""))
    page_break()

    # ── CHECKS AND BALANCES ───────────────────────────────────────────────────
    add_heading("4.  THEORY OF CHECKS AND BALANCES")
    add_divider()
    add_body(sections.get("checks_balances", ""))
    page_break()

    # ── INDIA CONTEXT ─────────────────────────────────────────────────────────
    add_heading("5.  SEPARATION OF POWERS IN INDIA")
    add_divider()
    add_body(sections.get("india_context", ""))
    page_break()

    # ── CONSTITUENT ASSEMBLY DEBATES ──────────────────────────────────────────
    add_heading("6.  THE CONSTITUENT ASSEMBLY DEBATES")
    add_divider()
    add_body(sections.get("assembly_debates", ""))
    page_break()

    # ── CONSTITUTIONAL REFLECTION ─────────────────────────────────────────────
    add_heading("7.  REFLECTION OF THE DOCTRINE IN THE INDIAN CONSTITUTION")
    add_divider()
    add_body(sections.get("constitutional_reflection", ""))
    page_break()

    # ── CASE LAW ──────────────────────────────────────────────────────────────
    add_heading("8.  LANDMARK CASE LAW AND JUDICIAL INTERPRETATION")
    add_divider()
    add_body(sections.get("case_law", ""))
    page_break()

    # ── CONCLUSION ───────────────────────────────────────────────────────────
    add_heading("9.  CONCLUSION")
    add_divider()
    add_body(sections.get("conclusion", ""))
    page_break()

    # ── BIBLIOGRAPHY ─────────────────────────────────────────────────────────
    add_heading("10. BIBLIOGRAPHY")
    p   = doc.add_paragraph()
    run = p.add_run("(Bluebook Citation Format)")
    set_font(run, 10, color=RGBColor(0x64, 0x64, 0x64))
    add_divider()
    add_body(sections.get("bibliography", ""))

    buf = _io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ── PDF builder ───────────────────────────────────────────────────────────────

def _assign_build_pdf(record: dict) -> bytes:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, PageBreak, HRFlowable, Table, TableStyle
        )
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
        import re as _re
    except ImportError:
        raise RuntimeError("reportlab not installed. Run: pip install reportlab")

    meta     = record["meta"]
    sections = record["sections"]
    buf      = _io.BytesIO()

    # CUSB palette
    CRIMSON = colors.HexColor("#CC0000")
    NAVY    = colors.HexColor("#1A3A5C")
    GOLD    = colors.HexColor("#C5922A")
    GREY    = colors.HexColor("#64748B")
    BLK     = colors.black

    # Double red border on every page
    def _draw_page_border(canvas_obj, doc_obj):
        w, h = A4
        canvas_obj.saveState()
        canvas_obj.setStrokeColor(CRIMSON)
        canvas_obj.setLineWidth(2.2)
        canvas_obj.rect(0.42*inch, 0.42*inch, w - 0.84*inch, h - 0.84*inch)
        canvas_obj.setLineWidth(0.7)
        canvas_obj.rect(0.52*inch, 0.52*inch, w - 1.04*inch, h - 1.04*inch)
        canvas_obj.restoreState()

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=1.05*inch, bottomMargin=0.8*inch,
        leftMargin=0.9*inch, rightMargin=0.9*inch,
        onFirstPage=_draw_page_border,
        onLaterPages=_draw_page_border,
    )

    S = {
        "cover_univ":   ParagraphStyle("cover_univ",   fontName="Helvetica-Bold",
                                        fontSize=18, textColor=CRIMSON,
                                        alignment=TA_CENTER, spaceAfter=4, leading=22),
        "cover_school": ParagraphStyle("cover_school", fontName="Helvetica-Bold",
                                        fontSize=12, textColor=CRIMSON,
                                        alignment=TA_CENTER, spaceAfter=6),
        "cover_sub_hd": ParagraphStyle("cover_sub_hd", fontName="Helvetica",
                                        fontSize=11, textColor=NAVY,
                                        alignment=TA_CENTER, spaceAfter=4),
        "cover_label":  ParagraphStyle("cover_label",  fontName="Helvetica",
                                        fontSize=10, textColor=NAVY,
                                        alignment=TA_CENTER, spaceAfter=3),
        "cover_topic":  ParagraphStyle("cover_topic",  fontName="Helvetica-Bold",
                                        fontSize=16, textColor=NAVY,
                                        alignment=TA_CENTER, spaceAfter=6, leading=22),
        "h1": ParagraphStyle("h1", fontName="Helvetica-Bold", fontSize=13,
                              textColor=CRIMSON, spaceBefore=18, spaceAfter=8),
        "h2": ParagraphStyle("h2", fontName="Helvetica-Bold", fontSize=11,
                              textColor=CRIMSON, spaceBefore=12, spaceAfter=6),
        "body": ParagraphStyle("body", fontName="Times-Roman", fontSize=11,
                                textColor=BLK, alignment=TA_JUSTIFY,
                                spaceAfter=7, leading=16),
        "toc":  ParagraphStyle("toc",  fontName="Times-Roman", fontSize=11,
                                textColor=BLK, spaceAfter=5),
        "sig":  ParagraphStyle("sig",  fontName="Helvetica", fontSize=10,
                                textColor=GREY, spaceAfter=4),
        "biblio": ParagraphStyle("bib", fontName="Times-Roman", fontSize=10,
                                  textColor=BLK, spaceAfter=5, leading=14),
    }

    def hr():
        return HRFlowable(width="100%", thickness=1.8, color=CRIMSON,
                          spaceAfter=10, spaceBefore=6)

    def thin_hr():
        return HRFlowable(width="100%", thickness=0.6, color=GOLD,
                          spaceAfter=6, spaceBefore=4)

    def h1(text):
        return Paragraph(f'<u>{text}</u>', S["h1"])

    def h2(text):
        return Paragraph(f'<u>{text}</u>', S["h2"])

    story = []

    # Cover page — CUSB style
    story += [Spacer(1, 0.5*inch)]
    story.append(Paragraph(meta["college"].upper(), S["cover_univ"]))
    story.append(hr())
    story.append(Paragraph(meta["course"], S["cover_school"]))
    story.append(thin_hr())
    story += [Spacer(1, 0.2*inch)]
    story.append(Paragraph("Subject: Legal Method", S["cover_sub_hd"]))
    story += [Spacer(1, 0.08*inch)]
    story.append(Paragraph("Assignment Topic", S["cover_label"]))
    story += [Spacer(1, 0.06*inch)]
    story.append(Paragraph(meta["topic"].upper(), S["cover_topic"]))
    story += [Spacer(1, 0.25*inch)]
    story.append(hr())
    story += [Spacer(1, 0.4*inch)]

    # Submitted to / Submitted by — 2-column table
    sub_lbl = ParagraphStyle("sl", fontName="Helvetica-Bold", fontSize=10,
                              textColor=CRIMSON, spaceAfter=4)
    sub_val = ParagraphStyle("sv", fontName="Times-Roman", fontSize=10,
                              textColor=BLK, spaceAfter=3, leading=14)
    col_to = [Paragraph("<b>SUBMITTED TO:</b>", sub_lbl),
              Paragraph("Course Faculty", sub_val),
              Paragraph("Department of Law and Governance", sub_val),
              Paragraph(meta["college"], sub_val)]
    col_by = [Paragraph("<b>SUBMITTED BY:</b>", sub_lbl),
              Paragraph(meta["student_name"], sub_val),
              Paragraph(meta["course"], sub_val),
              Paragraph(f'Roll No: {meta["roll_no"]}', sub_val),
              Paragraph("Department of Law and Governance", sub_val)]
    max_r = max(len(col_to), len(col_by))
    while len(col_to) < max_r: col_to.append(Paragraph("", sub_val))
    while len(col_by) < max_r: col_by.append(Paragraph("", sub_val))
    cover_tbl = Table([[col_to[i], col_by[i]] for i in range(max_r)],
                      colWidths=["48%", "48%"])
    cover_tbl.setStyle(TableStyle([
        ("VALIGN", (0,0), (-1,-1), "TOP"),
        ("LEFTPADDING",  (0,0), (-1,-1), 6),
        ("RIGHTPADDING", (0,0), (-1,-1), 6),
        ("TOPPADDING",   (0,0), (-1,-1), 2),
        ("BOTTOMPADDING",(0,0), (-1,-1), 2),
        ("LINEAFTER", (0,0), (0,-1), 1, CRIMSON),
    ]))
    story.append(cover_tbl)
    story.append(PageBreak())

    # TOC
    story.append(h1("TABLE OF CONTENTS"))
    story.append(hr())
    toc_items = [
        ("Declaration",                                    "2"),
        ("Acknowledgment",                                 "3") if sections.get("acknowledgment") else None,
        ("1.  Introduction",                               "4"),
        ("2.  The Concept and Theoretical Foundations",    "6"),
        ("3.  Structural Versus Functional Dimension",     "7"),
        ("4.  Theory of Checks and Balances",              "9"),
        ("5.  Separation of Powers in India",              "11"),
        ("6.  The Constituent Assembly Debates",           "12"),
        ("7.  Reflection of the Doctrine in Constitution", "13"),
        ("8.  Landmark Case Law",                          "14"),
        ("9.  Conclusion",                                 "16"),
        ("10. Bibliography",                               "17"),
    ]
    for item in toc_items:
        if item:
            dots = "." * max(1, 58 - len(item[0]))
            story.append(Paragraph(
                f'{item[0]} <font color="#A89E8E">{dots}</font> {item[1]}', S["toc"]
            ))
    story.append(PageBreak())

    def add_section(heading_text, text, bib=False):
        story.append(h1(heading_text))
        story.append(hr())
        if not text:
            return
        for line in text.split("\n"):
            line = line.strip()
            if not line:
                story.append(Spacer(1, 4))
                continue
            is_sub = bool(_re.match(r'^(\d+\.\d+|[IVX]+\.)\s', line))
            safe   = line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            if is_sub:
                story.append(h2(safe))
            else:
                sty = S["biblio"] if bib else S["body"]
                story.append(Paragraph(safe, sty))

    add_section("DECLARATION", sections.get("declaration", ""))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph(
        f'Signature: ________________&nbsp;&nbsp;&nbsp;&nbsp;'
        f'Date: {_dt.now().strftime("%d %B %Y")}', S["sig"]
    ))
    story.append(PageBreak())

    if sections.get("acknowledgment"):
        add_section("ACKNOWLEDGMENT", sections["acknowledgment"])
        story.append(PageBreak())

    add_section("1.  INTRODUCTION",                                        sections.get("introduction", ""))
    story.append(PageBreak())
    add_section("2.  THE CONCEPT AND THEORETICAL FOUNDATIONS",             sections.get("concept", ""))
    story.append(PageBreak())
    add_section("3.  STRUCTURAL VERSUS FUNCTIONAL DIMENSION",              sections.get("structural_functional", ""))
    story.append(PageBreak())
    add_section("4.  THEORY OF CHECKS AND BALANCES",                       sections.get("checks_balances", ""))
    story.append(PageBreak())
    add_section("5.  SEPARATION OF POWERS IN INDIA",                       sections.get("india_context", ""))
    story.append(PageBreak())
    add_section("6.  THE CONSTITUENT ASSEMBLY DEBATES",                    sections.get("assembly_debates", ""))
    story.append(PageBreak())
    add_section("7.  REFLECTION OF THE DOCTRINE IN THE INDIAN CONSTITUTION", sections.get("constitutional_reflection", ""))
    story.append(PageBreak())
    add_section("8.  LANDMARK CASE LAW AND JUDICIAL INTERPRETATION",       sections.get("case_law", ""))
    story.append(PageBreak())
    add_section("9.  CONCLUSION",                                          sections.get("conclusion", ""))
    story.append(PageBreak())

    story.append(h1("10. BIBLIOGRAPHY"))
    story.append(Paragraph("(Bluebook Citation Format)", S["sig"]))
    story.append(hr())
    for line in (sections.get("bibliography") or "").split("\n"):
        line = line.strip()
        if line:
            safe = line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            story.append(Paragraph(safe, S["biblio"]))

    doc.build(story)
    buf.seek(0)
    return buf.read()


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/api/assignment/generate", methods=["POST"])
@limiter.limit("5 per minute")
def assign_generate():
    data     = request.get_json(silent=True) or {}
    required = ["topic", "student_name", "roll_no", "course", "college"]
    missing  = [f for f in required if not str(data.get(f, "")).strip()]
    if missing:
        return jsonify({"error": f"Missing fields: {', '.join(missing)}"}), 400

    topic = str(data["topic"]).strip()
    if len(topic) > 300:
        return jsonify({"error": "Topic too long (max 300 chars)"}), 400

    # Optional reference authors: [{ "author": "...", "book": "..." }, ...]
    raw_authors = data.get("reference_authors", [])
    reference_authors = []
    if isinstance(raw_authors, list):
        for entry in raw_authors:
            if isinstance(entry, dict):
                author = str(entry.get("author", "")).strip()
                book   = str(entry.get("book", "")).strip()
                if author:
                    reference_authors.append({"author": author, "book": book})

    meta = {
        "topic":                topic,
        "student_name":         str(data["student_name"]).strip(),
        "roll_no":              str(data["roll_no"]).strip(),
        "course":               str(data["course"]).strip(),
        "college":              str(data["college"]).strip(),
        "include_acknowledgment": bool(data.get("include_acknowledgment", True)),
        "reference_authors":    reference_authors,
    }

    try:
        record = _assign_assemble(meta)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    return jsonify({
        "success":       True,
        "assignment_id": record["id"],
        "sections":      list(record["sections"].keys()),
        "cached":        record.get("created_at", 0) < time.time() - 2,
    })


@app.route("/api/assignment/<aid>/docx", methods=["GET"])
def assign_download_docx(aid):
    record = _ASSIGN_STORE.get(aid)
    if not record:
        return jsonify({"error": "Assignment not found or expired"}), 404
    try:
        buf = _assign_build_docx(record)
    except RuntimeError as e:
        return jsonify({"error": str(e)}), 500
    safe = re.sub(r'[^a-zA-Z0-9_-]', '_', record["meta"]["topic"][:40]).strip('_')
    return send_file(
        _io.BytesIO(buf),
        as_attachment=True,
        download_name=f"{safe}_Assignment.docx",
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@app.route("/api/assignment/<aid>/pdf", methods=["GET"])
def assign_download_pdf(aid):
    record = _ASSIGN_STORE.get(aid)
    if not record:
        return jsonify({"error": "Assignment not found or expired"}), 404
    try:
        buf = _assign_build_pdf(record)
    except RuntimeError as e:
        return jsonify({"error": str(e)}), 500
    safe = re.sub(r'[^a-zA-Z0-9_-]', '_', record["meta"]["topic"][:40]).strip('_')
    return send_file(
        _io.BytesIO(buf),
        as_attachment=True,
        download_name=f"{safe}_Assignment.pdf",
        mimetype="application/pdf",
    )


@app.route("/api/assignment/<aid>/status", methods=["GET"])
def assign_status(aid):
    record = _ASSIGN_STORE.get(aid)
    if not record:
        return jsonify({"error": "Not found"}), 404
    return jsonify({
        "id":         record["id"],
        "topic":      record["meta"]["topic"],
        "student":    record["meta"]["student_name"],
        "sections":   list(record["sections"].keys()),
        "created_at": record["created_at"],
    })


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    app.run(debug=False)